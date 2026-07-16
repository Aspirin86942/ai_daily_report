"""Office/PDF bounded 文档解析器。"""

from __future__ import annotations

from collections.abc import Mapping
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from ..workers.contracts import WorkerParsePayload

OFFICE_PARSER_BACKEND = "office_v1"
PDF_TEXT_PARSER_BACKEND = "pdf_text_v1"
NOT_PARSED_PARSER_BACKEND = "not_parsed"

DEFAULT_DOCUMENT_EXCERPT_MAX_CHARS = 6000
DEFAULT_EXCEL_MAX_SHEETS = 5
DEFAULT_EXCEL_MAX_ROWS = 50
DEFAULT_EXCEL_MAX_COLUMNS = 20
DEFAULT_DOCX_MAX_PARAGRAPHS = 200
DEFAULT_DOCX_MAX_TABLES = 20
DEFAULT_DOCX_TABLE_MAX_ROWS = 50
DEFAULT_DOCX_TABLE_MAX_COLS = 12
DEFAULT_PPTX_MAX_SLIDES = 50


@dataclass(frozen=True, slots=True)
class DocumentParserOptions:
    """Office/PDF parser 运行选项。"""

    office_parser_backend: str = OFFICE_PARSER_BACKEND
    pdf_parser_backend: str = PDF_TEXT_PARSER_BACKEND
    include_pptx_notes: bool = True


class _PreviewBuilder:
    """按字符预算累计 Markdown-like preview。"""

    def __init__(self, max_chars: int) -> None:
        self.max_chars = max(1, int(max_chars))
        self.parts: list[str] = []
        self.current_chars = 0
        self.was_truncated = False

    def add(self, text: str) -> bool:
        """追加文本；达到预算后返回 False 并标记截断。"""
        if not text:
            return True
        separator = "\n\n" if self.parts else ""
        candidate = separator + text
        remaining = self.max_chars - self.current_chars
        if remaining <= 0:
            self.was_truncated = True
            return False
        if len(candidate) > remaining:
            self.parts.append(candidate[:remaining])
            self.current_chars = self.max_chars
            self.was_truncated = True
            return False
        self.parts.append(candidate)
        self.current_chars += len(candidate)
        return True

    def build(self) -> str:
        """返回已累计的 preview 文本。"""
        return "".join(self.parts).strip()


def parse_document_file(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, Any],
    options: DocumentParserOptions | None = None,
) -> WorkerParsePayload:
    """按扩展名解析 Office/PDF 文件并返回 worker 内部结果。"""
    parser_options = options or DocumentParserOptions()
    normalized_type = file_type.lower()
    try:
        if normalized_type == ".docx":
            return _parse_docx(file_path, normalized_type, limits, parser_options)
        if normalized_type == ".xlsx":
            return _parse_xlsx(file_path, normalized_type, limits, parser_options)
        if normalized_type == ".pptx":
            return _parse_pptx(file_path, normalized_type, limits, parser_options)
        if normalized_type == ".pdf":
            return _parse_pdf(file_path, normalized_type, limits, parser_options)
        return WorkerParsePayload(
            file_path=str(file_path),
            file_type=normalized_type,
            content="",
            error=f"DOCUMENT_UNSUPPORTED_EXTENSION: {normalized_type}",
            parser_backend=NOT_PARSED_PARSER_BACKEND,
            truncated=False,
        )
    except Exception as exc:
        return WorkerParsePayload(
            file_path=str(file_path),
            file_type=normalized_type,
            content="",
            error=f"{_error_code_for(normalized_type)}: {exc}",
            parser_backend=_backend_for(normalized_type, parser_options),
            truncated=False,
        )


def _parse_docx(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, Any],
    options: DocumentParserOptions,
) -> WorkerParsePayload:
    from docx import Document

    document = Document(file_path)
    builder = _new_builder(limits)
    truncated = False

    paragraphs = [item.text.strip() for item in document.paragraphs if item.text.strip()]
    max_paragraphs = _positive_limit(
        limits,
        "docx_max_paragraphs",
        DEFAULT_DOCX_MAX_PARAGRAPHS,
    )
    if paragraphs:
        if not builder.add("# DOCX preview"):
            truncated = True
        if not builder.add("## Paragraphs"):
            truncated = True
        for index, paragraph in enumerate(paragraphs):
            if index >= max_paragraphs:
                truncated = True
                break
            if not builder.add(paragraph):
                truncated = True
                break
        if len(paragraphs) > max_paragraphs:
            truncated = True

    max_tables = _positive_limit(limits, "docx_max_tables", DEFAULT_DOCX_MAX_TABLES)
    max_rows = _positive_limit(
        limits,
        "docx_table_max_rows",
        DEFAULT_DOCX_TABLE_MAX_ROWS,
    )
    max_cols = _positive_limit(
        limits,
        "docx_table_max_cols",
        DEFAULT_DOCX_TABLE_MAX_COLS,
    )
    for table_index, table in enumerate(document.tables):
        if table_index >= max_tables:
            truncated = True
            break
        markdown, table_truncated = _docx_table_to_markdown(table, max_rows, max_cols)
        truncated = truncated or table_truncated
        if markdown:
            if not builder.add(f"## Table {table_index + 1}\n\n{markdown}"):
                truncated = True
                break

    content = builder.build()
    if not content:
        content = "No paragraph/table text extracted"
    return WorkerParsePayload(
        file_path=str(file_path),
        file_type=file_type,
        content=content,
        error=None,
        parser_backend=options.office_parser_backend,
        truncated=truncated or builder.was_truncated,
    )


def _parse_xlsx(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, Any],
    options: DocumentParserOptions,
) -> WorkerParsePayload:
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, read_only=True, data_only=True)
    builder = _new_builder(limits)
    truncated = False
    max_sheets = _positive_limit(limits, "excel_max_sheets", DEFAULT_EXCEL_MAX_SHEETS)
    max_rows = _positive_limit(limits, "excel_max_rows", DEFAULT_EXCEL_MAX_ROWS)
    max_cols = _positive_limit(limits, "excel_max_columns", DEFAULT_EXCEL_MAX_COLUMNS)

    try:
        if not builder.add("# XLSX preview"):
            truncated = True
        for sheet_index, sheet_name in enumerate(workbook.sheetnames):
            if sheet_index >= max_sheets:
                truncated = True
                break
            sheet = workbook[sheet_name]
            rows: list[list[str]] = []
            for raw_row in sheet.iter_rows(
                max_row=max_rows + 1,
                max_col=max_cols + 1,
                values_only=True,
            ):
                if _is_empty_row(raw_row):
                    continue
                if len(raw_row) > max_cols and any(
                    cell not in (None, "") for cell in raw_row[max_cols:]
                ):
                    truncated = True
                if len(rows) >= max_rows:
                    truncated = True
                    break
                rows.append([_format_cell(cell) for cell in raw_row[:max_cols]])
            markdown = _rows_to_markdown(rows)
            if markdown:
                if not builder.add(f"## Sheet: {sheet_name}\n\n{markdown}"):
                    truncated = True
                    break
    finally:
        workbook.close()

    content = builder.build() or "No worksheet text extracted"
    return WorkerParsePayload(
        file_path=str(file_path),
        file_type=file_type,
        content=content,
        error=None,
        parser_backend=options.office_parser_backend,
        truncated=truncated or builder.was_truncated,
    )


def _parse_pptx(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, Any],
    options: DocumentParserOptions,
) -> WorkerParsePayload:
    from pptx import Presentation

    presentation = Presentation(file_path)
    builder = _new_builder(limits)
    truncated = False
    max_slides = _positive_limit(limits, "pptx_max_slides", DEFAULT_PPTX_MAX_SLIDES)

    if not builder.add("# PPTX preview"):
        truncated = True
    for slide_index, slide in enumerate(presentation.slides):
        if slide_index >= max_slides:
            truncated = True
            break
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                markdown = _pptx_table_to_markdown(shape.table)
                if markdown:
                    slide_parts.append(markdown)
                continue
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_parts.append(text.strip())
        notes_text = _extract_pptx_notes(slide) if options.include_pptx_notes else ""
        if notes_text:
            slide_parts.append(f"### Notes\n\n{notes_text}")
        if slide_parts:
            if not builder.add(
                f"## Slide {slide_index + 1}\n\n" + "\n\n".join(slide_parts)
            ):
                truncated = True
                break

    content = builder.build() or "No slide text extracted"
    return WorkerParsePayload(
        file_path=str(file_path),
        file_type=file_type,
        content=content,
        error=None,
        parser_backend=options.office_parser_backend,
        truncated=truncated or builder.was_truncated,
    )


def _parse_pdf(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, Any],
    options: DocumentParserOptions,
) -> WorkerParsePayload:
    import pdfplumber

    builder = _new_builder(limits)
    truncated = False
    extracted_any_text = False
    max_pages = _positive_limit(limits, "pdf_max_pages", 5)

    if not builder.add("# PDF text preview"):
        truncated = True
    with pdfplumber.open(file_path) as pdf:
        for page_index, page in enumerate(pdf.pages):
            if page_index >= max_pages:
                truncated = True
                break
            text = page.extract_text() or ""
            if text.strip():
                extracted_any_text = True
                if not builder.add(f"## Page {page_index + 1}\n\n{text.strip()}"):
                    truncated = True
                    _close_pdf_page(page)
                    break
            _close_pdf_page(page)
        if len(pdf.pages) > max_pages:
            truncated = True

    if not extracted_any_text:
        return WorkerParsePayload(
            file_path=str(file_path),
            file_type=file_type,
            content="",
            error=f"PDF_NO_TEXT_LAYER: no extractable text in first {max_pages} pages",
            parser_backend=options.pdf_parser_backend,
            truncated=False,
        )

    return WorkerParsePayload(
        file_path=str(file_path),
        file_type=file_type,
        content=builder.build(),
        error=None,
        parser_backend=options.pdf_parser_backend,
        truncated=truncated or builder.was_truncated,
    )


def _new_builder(limits: Mapping[str, Any]) -> _PreviewBuilder:
    return _PreviewBuilder(
        _positive_limit(
            limits,
            "document_excerpt_max_chars",
            int(limits.get("text_max_chars", DEFAULT_DOCUMENT_EXCERPT_MAX_CHARS)),
        )
    )


def _positive_limit(
    limits: Mapping[str, Any],
    key: str,
    default: int,
) -> int:
    try:
        value = int(limits.get(key, default))
    except (TypeError, ValueError):
        return default
    return value if value > 0 else default


def _docx_table_to_markdown(table: Any, max_rows: int, max_cols: int) -> tuple[str, bool]:
    rows: list[list[str]] = []
    truncated = False
    for row_index, row in enumerate(table.rows):
        if row_index >= max_rows:
            truncated = True
            break
        cells = row.cells
        if len(cells) > max_cols:
            truncated = True
        rows.append([_clean_cell(cell.text) for cell in cells[:max_cols]])
    return _rows_to_markdown(rows), truncated


def _pptx_table_to_markdown(table: Any) -> str:
    rows = [
        [_clean_cell(cell.text) for cell in row.cells]
        for row in table.rows
    ]
    return _rows_to_markdown(rows)


def _rows_to_markdown(rows: list[list[str]]) -> str:
    normalized_rows = [row for row in rows if any(cell.strip() for cell in row)]
    if not normalized_rows:
        return ""
    column_count = max(len(row) for row in normalized_rows)
    padded_rows = [row + [""] * (column_count - len(row)) for row in normalized_rows]
    header = padded_rows[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in padded_rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _extract_pptx_notes(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
    except Exception:
        return ""
    if notes_frame is None:
        return ""
    return notes_frame.text.strip()


def _close_pdf_page(page: Any) -> None:
    close = getattr(page, "close", None)
    if callable(close):
        close()


def _is_empty_row(row: tuple[Any, ...]) -> bool:
    return all(cell is None or str(cell).strip() == "" for cell in row)


def _format_cell(value: Any) -> str:
    if value is None:
        return ""
    return _clean_cell(str(value))


def _clean_cell(value: str) -> str:
    # Markdown 表格用 pipe 分隔，替换单元格内 pipe 避免破坏列结构。
    return " ".join(value.replace("|", "/").split())


def _backend_for(file_type: str, options: DocumentParserOptions) -> str:
    if file_type == ".pdf":
        return options.pdf_parser_backend
    if file_type in {".docx", ".xlsx", ".pptx"}:
        return options.office_parser_backend
    return NOT_PARSED_PARSER_BACKEND


def _error_code_for(file_type: str) -> str:
    return {
        ".docx": "DOCX_PARSE_FAILED",
        ".xlsx": "XLSX_PARSE_FAILED",
        ".pptx": "PPTX_PARSE_FAILED",
        ".pdf": "PDF_PARSE_FAILED",
    }.get(file_type, "DOCUMENT_PARSE_FAILED")
