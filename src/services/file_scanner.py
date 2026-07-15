"""文件扫描服务"""

import multiprocessing as mp
from pathlib import Path
from datetime import date, timedelta
from typing import List, Mapping, Optional
from time import perf_counter
import pandas as pd
import pdfplumber
from pptx import Presentation

from ..models.schemas import FileContext, ScanResult
from ..core.config import config
from ..core.logger import setup_logger
from ..utils.text_tools import truncate_text
from .scan_discovery import DiscoveredFile, FileDiscoveryService
from .scan_index_store import InventoryItem, ScanIndexStore
from .light_text_parser import (
    DEFAULT_TEXT_MAX_CHARS,
    LIGHT_TEXT_PARSER_BACKEND,
    LightTextParserOptions,
    build_light_text_budget,
    parse_text_like_file,
)
from .document_parser import (
    DocumentParserOptions,
    parse_document_file,
)
from .office_parser import (
    OFFICE_RUST_FILE_TYPES,
    OfficeParseAudit,
    OfficeParseOutcome,
    PYTHON_OFFICE_BACKEND,
    PYTHON_SHAREPOINT_TEXT_BACKEND,
    RUST_OFFICE_BACKEND,
    parse_office_with_fallback,
    parse_with_sharepoint_text,
)
from .scan_metrics import ReparseDetail
from .scan_planner import DEFAULT_RUST_OFFICE_PARSER_BIN, ScanPlanner
from .scan_worker_pool import ParserSupervisor
from .rust_cli_contract import resolve_binary_path
from .cold_scanner_run import ColdScannerRun
from .scanner_items import (
    item_extension,
    item_identity,
    item_path,
    item_source_version,
    normalize_discovered_files,
)
from .scanner_parse_cache import (
    build_reparse_detail,
    build_reparse_exception_detail,
    get_cached_contexts,
    write_parse_cache,
)

logger = setup_logger()

TEXT_FILE_TYPES = {".txt", ".md", ".csv", ".json", ".log"}
DOCUMENT_FILE_TYPES = {".docx", ".xlsx", ".pptx", ".pdf"}
NOT_PARSED_PARSER_BACKEND = "not_parsed"
MODERN_OFFICE_FILE_TYPES = {".docx", ".xlsx", ".pptx"}
LEGACY_SHAREPOINT_FILE_TYPES = {".doc", ".ppt"}
DEFAULT_OFFICE_FALLBACK_ORDER = (
    PYTHON_OFFICE_BACKEND,
    PYTHON_SHAREPOINT_TEXT_BACKEND,
)


def _extract_content_worker(
    file_path_str: str,
    limits: dict,
    scanner_cfg: dict,
    result_queue: mp.Queue,
) -> None:
    """子进程解析单个文件并返回可序列化结果。"""
    scanner = object.__new__(FileScanner)
    scanner.scanner_cfg = scanner_cfg
    scanner.work_dir = Path(".")
    context = scanner._extract_content(Path(file_path_str), limits)
    result_queue.put(context.model_dump())


def _extract_document_content_worker(
    file_path_str: str,
    file_type: str,
    limits: dict,
    scanner_cfg: dict,
    result_queue: mp.Queue,
) -> None:
    """子进程解析 Office/PDF 文件并返回可序列化结果。"""
    context = parse_document_file(
        file_path=Path(file_path_str),
        file_type=file_type,
        limits=limits,
        options=DocumentParserOptions(
            office_parser_backend=scanner_cfg.get("office_parser_backend", "office_v1"),
            pdf_parser_backend=scanner_cfg.get("pdf_parser_backend", "pdf_text_v1"),
            include_pptx_notes=bool(scanner_cfg.get("pptx_include_notes", True)),
        ),
    )
    result_queue.put(context.model_dump())


def _extract_python_office_fallback_worker(
    file_path_str: str,
    file_type: str,
    limits: dict,
    scanner_cfg: dict,
    result_queue: mp.Queue,
) -> None:
    """子进程运行 Office Python fallback，避免坏文件卡住扫描线程。"""
    file_path = Path(file_path_str)
    normalized_type = file_type.lower()
    try:
        context = _run_python_office_fallback_backend(
            file_path=file_path,
            file_type=normalized_type,
            limits=limits,
            scanner_cfg=scanner_cfg,
        )
    except Exception as exc:
        context = FileContext(
            file_path=str(file_path),
            file_type=normalized_type,
            content="",
            error=f"PYTHON_OFFICE_FALLBACK_FAILED: {exc}",
            parser_backend=NOT_PARSED_PARSER_BACKEND,
            truncated=False,
        )
    result_queue.put(context.model_dump())


def _extract_configured_python_office_backend_worker(
    file_path_str: str,
    file_type: str,
    limits: dict,
    scanner_cfg: dict,
    backend: str,
    result_queue: mp.Queue,
) -> None:
    """子进程运行显式配置的 Python Office backend。"""
    file_path = Path(file_path_str)
    normalized_type = file_type.lower()
    try:
        context = _run_configured_python_office_backend(
            file_path=file_path,
            file_type=normalized_type,
            limits=limits,
            scanner_cfg=scanner_cfg,
            backend=backend,
        )
    except Exception as exc:
        context = FileContext(
            file_path=str(file_path),
            file_type=normalized_type,
            content="",
            error=f"PYTHON_OFFICE_BACKEND_FAILED: {exc}",
            parser_backend=NOT_PARSED_PARSER_BACKEND,
            truncated=False,
        )
    result_queue.put(context.model_dump())


def _run_python_office_fallback_backend(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, object],
    scanner_cfg: Mapping[str, object],
) -> FileContext:
    """按 scanner fallback_order 在子进程里选择 Python Office fallback。"""
    normalized_type = file_type.lower()
    last_context: FileContext | None = None

    for backend in _resolve_office_fallback_order(scanner_cfg):
        context: FileContext | None = None
        if backend == PYTHON_OFFICE_BACKEND:
            context = _run_python_office_backend(
                file_path=file_path,
                file_type=normalized_type,
                limits=limits,
                scanner_cfg=scanner_cfg,
            )
        elif (
            backend == PYTHON_SHAREPOINT_TEXT_BACKEND
            and (
                normalized_type in LEGACY_SHAREPOINT_FILE_TYPES
                or (
                    normalized_type == ".xls"
                    and last_context is not None
                    and last_context.parser_backend == PYTHON_OFFICE_BACKEND
                )
            )
        ):
            context = parse_with_sharepoint_text(
                file_path,
                normalized_type,
                limits,
            )

        if context is None:
            continue
        if context.error is None:
            return context
        last_context = context

    if last_context is not None:
        return last_context

    return FileContext(
        file_path=str(file_path),
        file_type=normalized_type,
        content="",
        error=f"PYTHON_FALLBACK_UNAVAILABLE: {normalized_type}",
        parser_backend=NOT_PARSED_PARSER_BACKEND,
        truncated=False,
    )


def _run_configured_python_office_backend(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, object],
    scanner_cfg: Mapping[str, object],
    backend: str,
) -> FileContext:
    """按 office_parser_backend 精确执行 Python backend。"""
    normalized_type = file_type.lower()
    if backend == PYTHON_OFFICE_BACKEND:
        context = _run_python_office_backend(
            file_path=file_path,
            file_type=normalized_type,
            limits=limits,
            scanner_cfg=scanner_cfg,
        )
        if context is not None:
            return context
        return FileContext(
            file_path=str(file_path),
            file_type=normalized_type,
            content="",
            error=f"PYTHON_OFFICE_UNSUPPORTED_EXTENSION: {normalized_type}",
            parser_backend=NOT_PARSED_PARSER_BACKEND,
            truncated=False,
        )
    if backend == PYTHON_SHAREPOINT_TEXT_BACKEND:
        return parse_with_sharepoint_text(file_path, normalized_type, limits)
    return FileContext(
        file_path=str(file_path),
        file_type=normalized_type,
        content="",
        error=f"OFFICE_UNKNOWN_BACKEND: {backend}",
        parser_backend=NOT_PARSED_PARSER_BACKEND,
        truncated=False,
    )


def _run_python_office_backend(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, object],
    scanner_cfg: Mapping[str, object],
) -> FileContext | None:
    """执行 python_office_v1 fallback；.xls 保留旧 Excel 表格抽取。"""
    if file_type in MODERN_OFFICE_FILE_TYPES:
        return parse_document_file(
            file_path=file_path,
            file_type=file_type,
            limits=dict(limits),
            options=_build_python_office_fallback_options(scanner_cfg),
        )
    if file_type == ".xls":
        return _parse_legacy_excel_fallback(
            file_path=file_path,
            file_type=file_type,
            limits=limits,
            scanner_cfg=scanner_cfg,
        )
    return None


def _parse_legacy_excel_fallback(
    file_path: Path,
    file_type: str,
    limits: Mapping[str, object],
    scanner_cfg: Mapping[str, object],
) -> FileContext:
    """旧 .xls fallback 继续输出表格 Markdown，保持历史 scanner 行为。"""
    max_rows = _positive_int(
        limits.get("excel_max_rows"),
        _positive_int(scanner_cfg.get("excel_max_rows"), 50),
    )
    text_max_chars = _positive_int(
        limits.get("text_max_chars"),
        _positive_int(scanner_cfg.get("text_max_chars"), 6000),
    )
    try:
        raw_content = _parse_excel_table_content(file_path, max_rows)
        truncated = len(raw_content) > text_max_chars
        content = truncate_text(raw_content, text_max_chars)
        return FileContext(
            file_path=str(file_path),
            file_type=file_type,
            content=content,
            error=None,
            parser_backend=PYTHON_OFFICE_BACKEND,
            truncated=truncated,
        )
    except Exception as exc:
        return FileContext(
            file_path=str(file_path),
            file_type=file_type,
            content="",
            error=f"PYTHON_OFFICE_XLS_FAILED: {exc}",
            parser_backend=PYTHON_OFFICE_BACKEND,
            truncated=False,
        )


def _parse_excel_table_content(file_path: Path, max_rows: int) -> str:
    """解析 Excel 文件为 Markdown 表格文本，供旧 scanner 和 .xls fallback 共用。"""
    content_parts = []
    excel_file = pd.ExcelFile(file_path)

    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=max_rows)
        df = df.dropna(how="all")

        if len(df) > max_rows:
            df = df.head(max_rows)
            content_parts.append(f"## {sheet_name} (仅显示前 {max_rows} 行)")
        else:
            content_parts.append(f"## {sheet_name}")

        if not df.empty:
            content_parts.append(df.to_markdown(index=False))

    return "\n\n".join(content_parts)


def _build_python_office_fallback_options(
    scanner_cfg: Mapping[str, object],
) -> DocumentParserOptions:
    """Python fallback 固定标记独立 backend，避免和旧 office_v1 混淆。"""
    return DocumentParserOptions(
        office_parser_backend=PYTHON_OFFICE_BACKEND,
        pdf_parser_backend=scanner_cfg.get(
            "pdf_parser_backend",
            "pdf_text_v1",
        ),
        include_pptx_notes=bool(scanner_cfg.get("pptx_include_notes", True)),
    )


def _resolve_office_fallback_order(
    scanner_cfg: Mapping[str, object],
) -> tuple[str, ...]:
    order = scanner_cfg.get("office_parser_fallback_order", DEFAULT_OFFICE_FALLBACK_ORDER)
    if isinstance(order, str):
        return (order,)
    try:
        return tuple(str(item) for item in order)
    except TypeError:
        return DEFAULT_OFFICE_FALLBACK_ORDER


def _positive_int(value: object, default: int) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        return default
    return parsed if parsed > 0 else default


class FileScanner:
    """文件扫描器"""

    def __init__(self):
        """初始化扫描器"""
        self.scanner_cfg = config.scanner_config
        self.work_dir = config.work_dir
        self.discovery_service = FileDiscoveryService(self.work_dir, self.scanner_cfg)
        (
            rust_office_binary_size_bytes,
            rust_office_binary_mtime_ns,
        ) = self._rust_office_binary_metadata(self.scanner_cfg)
        self.scan_planner = ScanPlanner(
            self.scanner_cfg,
            rust_office_parser_bin_size_bytes=rust_office_binary_size_bytes,
            rust_office_parser_bin_mtime_ns=rust_office_binary_mtime_ns,
        )
        self.scan_index_store = ScanIndexStore(
            self._resolve_project_path(self.scanner_cfg["index_db_path"])
        )
        self.parser_supervisor = ParserSupervisor(
            file_timeout_seconds=self.scanner_cfg.get("file_timeout_seconds", 30.0),
            file_timeout_by_extension=(
                self.scanner_cfg.get("file_timeout_by_extension", {}) or {}
            ),
        )
        self.last_reparse_details: list[ReparseDetail] = []
        self._office_parse_audits: dict[str, OfficeParseAudit] = {}

    @staticmethod
    def _resolve_project_path(path_value: str | Path) -> Path:
        """把相对配置路径解析到项目根目录，避免随运行目录漂移。"""
        path = Path(path_value)
        if path.is_absolute():
            return path
        project_root = Path(__file__).resolve().parent.parent.parent
        return project_root / path

    @staticmethod
    def _rust_office_binary_metadata(
        scanner_cfg: dict,
    ) -> tuple[int | None, int | None]:
        """在 runtime adapter 解析 helper 指纹，planner 保持无文件系统 I/O。"""
        binary_path = resolve_binary_path(
            scanner_cfg.get(
                "rust_office_parser_bin",
                DEFAULT_RUST_OFFICE_PARSER_BIN,
            )
        )
        try:
            binary_stat = binary_path.stat()
        except OSError:
            return None, None
        return binary_stat.st_size, binary_stat.st_mtime_ns

    def scan_today_files(self) -> ScanResult:
        """扫描今日修改的文件（默认日期范围封装）

        Returns:
            扫描汇总结果
        """
        today = date.today()
        yesterday = today - timedelta(days=1)
        return self.scan_files(start_date=yesterday, end_date=today)

    def scan_files(
        self,
        start_date: Optional[date] = None,
        end_date: Optional[date] = None,
        summary_mode: bool = False,
    ) -> ScanResult:
        """通用文件扫描方法

        Args:
            start_date: 起始日期 (默认昨日)
            end_date: 结束日期 (默认今日)
            summary_mode: 摘要模式 — 使用缩减的解析限制

        Returns:
            扫描汇总结果
        """
        return ColdScannerRun(self).scan_files(
            start_date=start_date,
            end_date=end_date,
            summary_mode=summary_mode,
        )

    def _normalize_discovered_files(
        self,
        discovered_files: list[Path | DiscoveredFile],
    ) -> list[DiscoveredFile]:
        """兼容旧 Path monkeypatch，同时统一生成 inventory 所需元数据。"""
        return normalize_discovered_files(discovered_files)

    def _get_cached_contexts(
        self,
        cached_files: list[Path | InventoryItem],
        parser_profile: str,
    ) -> list[FileContext]:
        """从 parse_cache 恢复 fresh cache 命中的上下文。"""
        return get_cached_contexts(
            self.scan_index_store,
            cached_files,
            parser_profile,
        )

    def _get_files_in_range(self, start_date: date, end_date: date) -> List[Path]:
        """获取日期范围内修改的文件列表

        Args:
            start_date: 起始日期
            end_date: 结束日期

        Returns:
            文件路径列表
        """
        discovered_files = self.discovery_service.bootstrap_full_scan(start_date, end_date)
        return [
            item.path if isinstance(item, DiscoveredFile) else Path(item)
            for item in discovered_files
        ]

    def _item_path(self, item: Path | InventoryItem) -> Path:
        """统一读取候选路径。"""
        return item_path(item)

    def _item_identity(self, item: Path | InventoryItem) -> str:
        """统一读取缓存身份。"""
        return item_identity(item)

    def _item_extension(self, item: Path | InventoryItem) -> str:
        """统一读取扩展名。"""
        return item_extension(item)

    def _item_source_version(self, item: Path | InventoryItem) -> str:
        """统一读取 discovery 版本指纹。"""
        return item_source_version(item)

    def _write_parse_cache(
        self,
        item: Path | InventoryItem,
        parser_profile: str,
        context: FileContext,
    ) -> None:
        """把本轮解析结果写回 parse_cache。"""
        write_parse_cache(self.scan_index_store, item, parser_profile, context)

    def _record_reparse_detail(
        self,
        item: Path | InventoryItem,
        cache_probe,
        duration_ms: int,
        context: FileContext,
    ) -> None:
        """记录单个重解析文件的 cache miss 原因和解析结果。"""
        self.last_reparse_details.append(
            build_reparse_detail(
                item=item,
                cache_probe=cache_probe,
                duration_ms=duration_ms,
                context=context,
                office_parse_audits=self._office_parse_audits,
                infer_worker_lane=self._infer_worker_lane,
            )
        )

    def _record_reparse_exception(
        self,
        item: Path | InventoryItem,
        cache_probe,
        parse_error: str,
    ) -> None:
        """解析入口抛异常时，也要留下 benchmark 可见的重解析明细。"""
        self.last_reparse_details.append(
            build_reparse_exception_detail(
                item=item,
                cache_probe=cache_probe,
                parse_error=parse_error,
                not_parsed_backend=NOT_PARSED_PARSER_BACKEND,
            )
        )

    def _extract_content_with_timeout(
        self, file_path: Path, limits: Optional[dict] = None
    ) -> FileContext:
        """带单文件时间预算的内容提取入口。"""
        file_type = file_path.suffix.lower()
        timeout_seconds = self.parser_supervisor.resolve_timeout(file_type)
        context, timed_out = self._run_extract_subprocess(
            file_path,
            limits,
            timeout_seconds,
        )
        if timed_out:
            logger.warning(
                "解析文件超时: %s (%ss)",
                file_path,
                f"{timeout_seconds:g}",
            )
            return self.parser_supervisor.handle_worker_timeout(file_path, file_type)
        if context is None:
            return self.parser_supervisor.handle_missing_result(file_path, file_type)
        return context

    def _extract_content_with_duration(
        self,
        file_path: Path,
        limits: Optional[dict] = None,
    ) -> tuple[FileContext, int]:
        """解析单文件并返回本 worker 内部 wall clock 耗时。"""
        started_at = perf_counter()
        context = self._extract_content_with_timeout(file_path, limits)
        duration_ms = int(round((perf_counter() - started_at) * 1000))
        return context, max(0, duration_ms)

    def _extract_uncached_content_with_duration(
        self,
        item: Path | InventoryItem,
        limits: Optional[dict] = None,
    ) -> tuple[FileContext, int]:
        """解析未缓存文件，并返回本 worker 内部 wall clock 耗时。"""
        started_at = perf_counter()
        context = self._extract_uncached_content(
            self._item_path(item),
            self._item_extension(item),
            limits,
        )
        duration_ms = int(round((perf_counter() - started_at) * 1000))
        return context, max(0, duration_ms)

    def _extract_uncached_content(
        self,
        file_path: Path,
        file_type: str,
        limits: Optional[dict] = None,
    ) -> FileContext:
        """根据文件类型选择 light text parser 或 subprocess timeout lane。"""
        effective_limits = limits or {}
        too_large_context = self._build_file_too_large_context(
            file_path,
            file_type,
        )
        if too_large_context is not None:
            return too_large_context

        if self._should_parse_direct(file_type):
            return parse_text_like_file(
                file_path=file_path,
                file_type=file_type,
                limits=effective_limits,
                options=self._build_light_text_options(effective_limits),
            )
        if self._should_parse_office_rust(file_type):
            return self._extract_office_content_with_timeout(
                file_path,
                file_type,
                effective_limits,
            )
        if self._should_parse_document_direct(file_type):
            return self._extract_document_content_with_timeout(
                file_path,
                effective_limits,
            )
        return self._extract_content_with_timeout(file_path, effective_limits)

    def _should_parse_direct(self, file_type: str) -> bool:
        """text-like 文件使用 bounded direct parser，避免 Windows spawn 固定开销。"""
        if str(self.scanner_cfg.get("worker_lane_mode", "direct")).lower() != "direct":
            return False
        return file_type.lower() in TEXT_FILE_TYPES

    def _should_parse_office_rust(self, file_type: str) -> bool:
        """direct lane 下 Office 文件交给 Rust Office parser orchestration。"""
        if str(self.scanner_cfg.get("worker_lane_mode", "direct")).lower() != "direct":
            return False
        return file_type.lower() in OFFICE_RUST_FILE_TYPES

    def _should_parse_document_direct(self, file_type: str) -> bool:
        """Office/PDF 使用正式 backend，但仍通过子进程保留 hard timeout。"""
        if str(self.scanner_cfg.get("worker_lane_mode", "direct")).lower() != "direct":
            return False
        return file_type.lower() in DOCUMENT_FILE_TYPES

    def _build_light_text_options(self, limits: dict) -> LightTextParserOptions:
        """把 scanner 配置转换为 light parser 的有界读取选项。"""
        light_text_budget = build_light_text_budget(
            self.scanner_cfg,
            text_max_chars=limits.get(
                "text_max_chars",
                self.scanner_cfg.get("text_max_chars", DEFAULT_TEXT_MAX_CHARS),
            ),
            default_text_max_chars=DEFAULT_TEXT_MAX_CHARS,
            on_invalid=self._warn_invalid_light_text_budget,
        )
        return LightTextParserOptions(
            read_head_bytes=light_text_budget.direct_text_read_bytes,
            read_tail_bytes=light_text_budget.log_tail_read_bytes,
            max_output_chars=light_text_budget.text_excerpt_max_chars,
            parser_backend_version=LIGHT_TEXT_PARSER_BACKEND,
        )

    def _build_document_parser_options(self) -> DocumentParserOptions:
        """把 scanner 配置转换为 document parser 选项。"""
        return DocumentParserOptions(
            office_parser_backend=self.scanner_cfg.get(
                "office_parser_backend",
                "office_v1",
            ),
            pdf_parser_backend=self.scanner_cfg.get(
                "pdf_parser_backend",
                "pdf_text_v1",
            ),
            include_pptx_notes=bool(self.scanner_cfg.get("pptx_include_notes", True)),
        )

    def _build_python_office_fallback_options(self) -> DocumentParserOptions:
        """Python fallback 固定标记独立 backend，避免和旧 office_v1 混淆。"""
        return _build_python_office_fallback_options(self.scanner_cfg)

    def _warn_invalid_light_text_budget(
        self,
        key: str,
        raw_value: object,
        default: int,
        reason: str,
    ) -> None:
        """运行时记录无效预算配置；planner 只负责静默归一化 cache key。"""
        if reason == "invalid":
            logger.warning(
                "%s 配置无效，使用默认值 %s: %r",
                key,
                default,
                raw_value,
            )
            return
        logger.warning(
            "%s 配置必须为正整数，使用默认值 %s: %r",
            key,
            default,
            raw_value,
        )

    def _build_file_too_large_context(
        self,
        file_path: Path,
        file_type: str,
    ) -> FileContext | None:
        """复用 scanner 既有 max_file_size_mb 策略，避免 direct lane 绕过门禁。"""
        max_file_size_mb = self.scanner_cfg.get("max_file_size_mb")
        if max_file_size_mb is None:
            return None

        max_bytes = float(max_file_size_mb) * 1024 * 1024
        file_size = file_path.stat().st_size
        if file_size <= max_bytes:
            return None

        return FileContext(
            file_path=str(file_path),
            file_type=file_type,
            content="",
            error=(
                f"file too large: {file_size} bytes exceeds "
                f"{max_file_size_mb} MB limit"
            ),
            parser_backend=NOT_PARSED_PARSER_BACKEND,
        )

    def _run_extract_subprocess(
        self,
        file_path: Path,
        limits: Optional[dict],
        timeout_seconds: float,
    ) -> tuple[Optional[FileContext], bool]:
        """在独立子进程中解析文件，返回结果和是否超时。"""
        ctx = mp.get_context("spawn")
        result_queue: mp.Queue = ctx.Queue(maxsize=1)
        process = ctx.Process(
            target=_extract_content_worker,
            args=(str(file_path), limits or {}, dict(self.scanner_cfg), result_queue),
        )
        process.start()
        process.join(timeout_seconds)

        if process.is_alive():
            process.terminate()
            process.join()
            return None, True

        try:
            payload = result_queue.get_nowait()
        except Exception:
            return None, False

        try:
            return FileContext(**payload), False
        except Exception as exc:
            logger.warning("子进程返回无效结果 %s: %s", file_path, exc)
            return (
                self.parser_supervisor.handle_invalid_payload(
                    file_path,
                    file_path.suffix.lower(),
                ),
                False,
            )

    def _extract_document_content_with_timeout(
        self,
        file_path: Path,
        limits: Optional[dict] = None,
    ) -> FileContext:
        """带 timeout 的 Office/PDF 正式 backend 入口。"""
        file_type = file_path.suffix.lower()
        timeout_seconds = self.parser_supervisor.resolve_timeout(file_type)
        context, timed_out = self._run_extract_document_subprocess(
            file_path,
            limits,
            timeout_seconds,
        )
        if timed_out:
            logger.warning(
                "解析文件超时: %s (%ss)",
                file_path,
                f"{timeout_seconds:g}",
            )
            return self.parser_supervisor.handle_worker_timeout(file_path, file_type)
        if context is None:
            return self.parser_supervisor.handle_missing_result(file_path, file_type)
        return context

    def _extract_office_content_with_timeout(
        self,
        file_path: Path,
        file_type: str,
        limits: Mapping[str, object],
    ) -> FileContext:
        """运行 Rust Office parser，并记录 Rust/Python fallback 审计信息。"""
        normalized_type = file_type.lower()
        timeout_seconds = self.parser_supervisor.resolve_timeout(normalized_type)
        configured_backend = str(
            self.scanner_cfg.get("office_parser_backend", RUST_OFFICE_BACKEND)
        )
        if configured_backend != RUST_OFFICE_BACKEND:
            context = self._parse_configured_python_office_backend(
                file_path,
                normalized_type,
                limits,
                configured_backend,
            )
            self._office_parse_audits[str(file_path)] = OfficeParseAudit(
                attempted_backend=configured_backend,
            )
            return context

        outcome: OfficeParseOutcome = parse_office_with_fallback(
            file_path=file_path,
            file_type=normalized_type,
            limits=limits,
            scanner_cfg=self.scanner_cfg,
            timeout_seconds=timeout_seconds,
            python_fallback=self._parse_python_office_fallback,
        )
        self._office_parse_audits[str(file_path)] = outcome.audit
        return outcome.context

    def _parse_configured_python_office_backend(
        self,
        file_path: Path,
        file_type: str,
        limits: Mapping[str, object],
        backend: str,
    ) -> FileContext:
        """显式 Python backend 也必须通过子进程保留 hard timeout。"""
        timeout_seconds = self.parser_supervisor.resolve_timeout(file_type)
        context, timed_out = self._run_configured_python_office_backend_subprocess(
            file_path,
            file_type,
            dict(limits),
            backend,
            timeout_seconds,
        )
        if timed_out:
            logger.warning(
                "Office Python backend 超时: %s (%ss)",
                file_path,
                f"{timeout_seconds:g}",
            )
            return self.parser_supervisor.handle_worker_timeout(file_path, file_type)
        if context is None:
            return self.parser_supervisor.handle_missing_result(file_path, file_type)
        return context

    def _parse_python_office_fallback(
        self,
        file_path: Path,
        file_type: str,
        limits: Mapping[str, object],
    ) -> FileContext:
        """Rust Office 失败后的 Python fallback，放入子进程保留 hard timeout。"""
        normalized_type = file_type.lower()
        timeout_seconds = self.parser_supervisor.resolve_timeout(normalized_type)
        context, timed_out = self._run_python_office_fallback_subprocess(
            file_path,
            normalized_type,
            dict(limits),
            timeout_seconds,
        )
        if timed_out:
            logger.warning(
                "Office Python fallback 超时: %s (%ss)",
                file_path,
                f"{timeout_seconds:g}",
            )
            return self.parser_supervisor.handle_worker_timeout(
                file_path,
                normalized_type,
            )
        if context is None:
            return self.parser_supervisor.handle_missing_result(
                file_path,
                normalized_type,
            )
        return context

    def _run_python_office_fallback_subprocess(
        self,
        file_path: Path,
        file_type: str,
        limits: dict,
        timeout_seconds: float,
    ) -> tuple[Optional[FileContext], bool]:
        """在独立子进程中运行 Python Office fallback。"""
        ctx = mp.get_context("spawn")
        result_queue: mp.Queue = ctx.Queue(maxsize=1)
        process = ctx.Process(
            target=_extract_python_office_fallback_worker,
            args=(
                str(file_path),
                file_type,
                limits,
                dict(self.scanner_cfg),
                result_queue,
            ),
        )
        process.start()
        process.join(timeout_seconds)

        if process.is_alive():
            process.terminate()
            process.join()
            return None, True

        try:
            payload = result_queue.get_nowait()
        except Exception:
            return None, False

        try:
            return FileContext(**payload), False
        except Exception as exc:
            logger.warning("Office Python fallback 返回无效结果 %s: %s", file_path, exc)
            return (
                self.parser_supervisor.handle_invalid_payload(
                    file_path,
                    file_type,
                ),
                False,
            )

    def _run_configured_python_office_backend_subprocess(
        self,
        file_path: Path,
        file_type: str,
        limits: dict,
        backend: str,
        timeout_seconds: float,
    ) -> tuple[Optional[FileContext], bool]:
        """在独立子进程中运行显式 Python Office backend。"""
        ctx = mp.get_context("spawn")
        result_queue: mp.Queue = ctx.Queue(maxsize=1)
        process = ctx.Process(
            target=_extract_configured_python_office_backend_worker,
            args=(
                str(file_path),
                file_type,
                limits,
                dict(self.scanner_cfg),
                backend,
                result_queue,
            ),
        )
        process.start()
        process.join(timeout_seconds)

        if process.is_alive():
            process.terminate()
            process.join()
            return None, True

        try:
            payload = result_queue.get_nowait()
        except Exception:
            return None, False

        try:
            return FileContext(**payload), False
        except Exception as exc:
            logger.warning("Office Python backend 返回无效结果 %s: %s", file_path, exc)
            return (
                self.parser_supervisor.handle_invalid_payload(
                    file_path,
                    file_type,
                ),
                False,
            )

    def _run_extract_document_subprocess(
        self,
        file_path: Path,
        limits: Optional[dict],
        timeout_seconds: float,
    ) -> tuple[Optional[FileContext], bool]:
        """在独立子进程中运行 Office/PDF backend，避免坏文件拖死主进程。"""
        ctx = mp.get_context("spawn")
        result_queue: mp.Queue = ctx.Queue(maxsize=1)
        process = ctx.Process(
            target=_extract_document_content_worker,
            args=(
                str(file_path),
                file_path.suffix.lower(),
                limits or {},
                dict(self.scanner_cfg),
                result_queue,
            ),
        )
        process.start()
        process.join(timeout_seconds)

        if process.is_alive():
            process.terminate()
            process.join()
            return None, True

        try:
            payload = result_queue.get_nowait()
        except Exception:
            return None, False

        try:
            return FileContext(**payload), False
        except Exception as exc:
            logger.warning("子进程返回无效结果 %s: %s", file_path, exc)
            return (
                self.parser_supervisor.handle_invalid_payload(
                    file_path,
                    file_path.suffix.lower(),
                ),
                False,
            )

    def _resolve_file_timeout(self, file_type: str) -> float:
        """兼容旧调用点，内部转发给 supervisor。"""
        return self.parser_supervisor.resolve_timeout(file_type)

    def load_discovery_checkpoint(self, discovery_key: str) -> str | None:
        """读取 discovery checkpoint 占位值，供后续增量发现接线。"""
        return self.scan_index_store.load_checkpoint(discovery_key)

    def save_discovery_checkpoint(
        self,
        discovery_key: str,
        checkpoint_value: str,
    ) -> None:
        """写入 discovery checkpoint 占位值，但当前扫描流程仍不依赖它。"""
        self.scan_index_store.save_checkpoint(discovery_key, checkpoint_value)

    def _extract_content(
        self, file_path: Path, limits: Optional[dict] = None
    ) -> FileContext:
        """提取文件内容

        Args:
            file_path: 文件路径
            limits: 解析限制参数 (excel_max_rows, pdf_max_pages, text_max_chars)

        Returns:
            文件上下文
        """
        if limits is None:
            limits = {
                "excel_max_rows": self.scanner_cfg["excel_max_rows"],
                "pdf_max_pages": self.scanner_cfg["pdf_max_pages"],
                "text_max_chars": self.scanner_cfg["text_max_chars"],
            }

        file_type = file_path.suffix.lower()

        try:
            too_large_context = self._build_file_too_large_context(
                file_path,
                file_type,
            )
            if too_large_context is not None:
                return too_large_context

            if file_type in DOCUMENT_FILE_TYPES:
                parser_options = (
                    self._build_python_office_fallback_options()
                    if file_type in MODERN_OFFICE_FILE_TYPES
                    else self._build_document_parser_options()
                )
                return parse_document_file(
                    file_path=file_path,
                    file_type=file_type,
                    limits=limits,
                    options=parser_options,
                )
            if file_type in [".xls"]:
                content = self._parse_excel(file_path, limits["excel_max_rows"])
            elif file_type in TEXT_FILE_TYPES:
                content = self._parse_text(file_path, limits["text_max_chars"])
            else:
                content = ""

            # 截断过长文本
            content = truncate_text(content, limits["text_max_chars"])

            return FileContext(
                file_path=str(file_path),
                file_type=file_type,
                content=content,
                error=None,
            )

        except Exception as e:
            logger.warning(f"解析文件失败 {file_path}: {e}")
            return FileContext(
                file_path=str(file_path), file_type=file_type, content="", error=str(e)
            )

    def _infer_worker_lane(self, file_type: str, context: FileContext) -> str:
        """区分执行通道和 parser backend，避免 benchmark 把两者混在一起。"""
        if context.parser_backend == NOT_PARSED_PARSER_BACKEND:
            return "not_parsed"
        if str(self.scanner_cfg.get("worker_lane_mode", "direct")).lower() != "direct":
            return "subprocess"
        if file_type.lower() in TEXT_FILE_TYPES:
            return "direct"
        if file_type.lower() in OFFICE_RUST_FILE_TYPES:
            return "subprocess"
        if file_type.lower() in DOCUMENT_FILE_TYPES:
            # Office/PDF 虽然是正式 backend，但仍由子进程提供 hard timeout 隔离。
            return "subprocess"
        return "subprocess"

    def _parse_excel(self, file_path: Path, max_rows: Optional[int] = None) -> str:
        """解析 Excel 文件

        Args:
            file_path: 文件路径
            max_rows: 最大行数限制

        Returns:
            Markdown 格式的表格内容
        """
        if max_rows is None:
            max_rows = self.scanner_cfg["excel_max_rows"]

        return _parse_excel_table_content(file_path, max_rows)

    def _parse_pdf(self, file_path: Path, max_pages: Optional[int] = None) -> str:
        """解析 PDF 文件

        Args:
            file_path: 文件路径
            max_pages: 最大页数限制

        Returns:
            提取的文本内容
        """
        if max_pages is None:
            max_pages = self.scanner_cfg["pdf_max_pages"]

        content_parts = []

        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages[:max_pages]):
                text = page.extract_text()
                if text:
                    content_parts.append(f"## 第 {i + 1} 页\n{text}")

            if len(pdf.pages) > max_pages:
                content_parts.append(
                    f"\n(PDF 共 {len(pdf.pages)} 页，仅显示前 {max_pages} 页)"
                )

        return "\n\n".join(content_parts)

    def _parse_pptx(self, file_path: Path) -> str:
        """解析 PPTX 文件

        Args:
            file_path: 文件路径

        Returns:
            提取的文本内容
        """
        content_parts = []
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            if slide_text:
                content_parts.append(f"## 幻灯片 {i + 1}\n" + "\n".join(slide_text))

        return "\n\n".join(content_parts)

    def _parse_text(self, file_path: Path, max_chars: Optional[int] = None) -> str:
        """解析纯文本文件 (.txt, .md)

        Args:
            file_path: 文件路径
            max_chars: 最大读取字符数

        Returns:
            文件文本内容
        """
        if max_chars is None:
            max_chars = self.scanner_cfg["text_max_chars"]
        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read(max_chars + 1)
        return content

    def _parse_docx(self, file_path: Path) -> str:
        """解析 Word 文档

        Args:
            file_path: 文件路径

        Returns:
            提取的文本内容
        """
        from docx import Document

        doc = Document(file_path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
