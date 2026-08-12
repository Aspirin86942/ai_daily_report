"""Worker v1 的无请求体 version 进程合同。"""

from __future__ import annotations

import json
from pathlib import Path
import subprocess
import sys

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from src.models.scanner_contract import (
    WorkerParseResponse,
)


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OFFICE_WORKER_BIN = (
    PROJECT_ROOT / "rust" / "target" / "release" / "ai-daily-office-parser.exe"
)


def _require_office_worker() -> None:
    if OFFICE_WORKER_BIN.is_file():
        return
    if sys.platform == "win32":
        pytest.fail(
            "Windows integration requires cargo build --manifest-path "
            "rust/Cargo.toml --workspace --release --locked"
        )
    pytest.skip("Rust Office worker release binary is not built")


def test_python_document_worker_hello_is_strict_requestless_json() -> None:
    completed = subprocess.run(
        [
            sys.executable,
            "-m",
            "src.workers.document_parser_worker",
            "hello",
        ],
        cwd=PROJECT_ROOT,
        input=b"",
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr.decode(
        "utf-8",
        errors="replace",
    )
    assert completed.stderr == b""
    payload = json.loads(completed.stdout.decode("utf-8", errors="strict"))
    assert payload["worker_kind"] == "python_document"
    assert payload["worker_contract_version"] == "ai_daily_worker_v2"
    assert payload["supported_operations"] == [
        "pdf_classify",
        "pdf_parse",
        "python_office_parse",
        "python_sharepoint_parse",
    ]


def test_python_document_worker_version_import_path_stays_stdlib_light() -> None:
    completed = subprocess.run(
        [
            sys.executable,
            "-S",
            "-c",
            (
                "import sys; "
                "import src.workers.document_parser_worker; "
                "forbidden={'json','pathlib','typing'} & set(sys.modules); "
                "raise SystemExit(1 if forbidden else 0)"
            ),
        ],
        cwd=PROJECT_ROOT,
        input=b"",
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr.decode(
        "utf-8",
        errors="replace",
    )


@pytest.mark.parametrize(
    ("file_name", "file_type", "backend", "expected_text"),
    [
        (
            "legacy_sample.xls",
            ".xls",
            "python_office_v2",
            "Legacy XLS worker content",
        ),
        (
            "legacy_sample.doc",
            ".doc",
            "python_sharepoint_text_v2",
            "Legacy DOC worker content",
        ),
        (
            "legacy_sample.ppt",
            ".ppt",
            "python_sharepoint_text_v2",
            "Legacy PPT worker content",
        ),
    ],
)
def test_python_document_worker_process_parses_real_legacy_office(
    file_name: str,
    file_type: str,
    backend: str,
    expected_text: str,
) -> None:
    sample = PROJECT_ROOT / "tests" / "fixtures" / "worker_documents" / file_name
    request = _python_document_parse_request(sample, file_type, backend)

    operation = (
        "python_office_parse"
        if backend == "python_office_v2"
        else "python_sharepoint_parse"
    )
    envelope = {
        "contract": "ai_daily_worker",
        "protocol_version": 2,
        "request_id": request["request_id"],
        "operation": operation,
        "payload": request,
    }
    completed = subprocess.run(
        [
            sys.executable,
            "-m",
            "src.workers.document_parser_worker",
            "session",
        ],
        cwd=PROJECT_ROOT,
        input=(json.dumps(envelope, ensure_ascii=False) + "\n").encode("utf-8"),
        capture_output=True,
        check=False,
        timeout=30,
    )

    assert completed.returncode == 0, completed.stderr.decode(
        "utf-8",
        errors="replace",
    )
    assert completed.stderr == b""
    frames = completed.stdout.splitlines()
    assert json.loads(frames[0])["frame"] == "hello"
    response = WorkerParseResponse.model_validate(json.loads(frames[1])["result"])
    assert response.status == "ok"
    assert response.parser_backend == backend
    assert response.worker_lane == "python_document_process_v2"
    assert expected_text in response.content
    assert response.observed_source_version == request["expected_source_version"]


def test_office_worker_hello_is_requestless_and_ignores_stdin() -> None:
    _require_office_worker()

    completed = subprocess.run(
        [str(OFFICE_WORKER_BIN), "hello"],
        cwd=PROJECT_ROOT,
        input=b"version commands do not read stdin",
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr.decode(
        "utf-8",
        errors="replace",
    )
    assert completed.stderr == b""
    payload = json.loads(completed.stdout.decode("utf-8", errors="strict"))
    assert payload["worker_kind"] == "office"
    assert payload["worker_contract_version"] == "ai_daily_worker_v2"
    assert payload["supported_operations"] == ["office_parse"]


def test_office_worker_requires_an_explicit_command() -> None:
    _require_office_worker()

    completed = subprocess.run(
        [str(OFFICE_WORKER_BIN)],
        cwd=PROJECT_ROOT,
        input=b"{}",
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 1
    assert completed.stdout == b""
    assert b"usage: ai-daily-office-parser <hello|session>" in completed.stderr


@pytest.mark.parametrize("file_type", [".xlsx", ".docx", ".pptx"])
def test_office_worker_strict_parse_handles_modern_office(
    tmp_path: Path,
    file_type: str,
) -> None:
    _require_office_worker()
    sample = tmp_path / f"strict Office{file_type}"
    if file_type == ".xlsx":
        workbook = Workbook()
        workbook.active.append(["Rust Office strict content"])
        workbook.save(sample)
        workbook.close()
        backend = "rust_xlsx_bounded_v2"
    elif file_type == ".docx":
        document = Document()
        document.add_paragraph("Rust Office strict content")
        document.save(sample)
        backend = "rust_office_oxide_v2"
    else:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Rust Office strict content"
        presentation.save(sample)
        backend = "rust_office_oxide_v2"
    request = _office_parse_request(sample, file_type, backend)

    completed = subprocess.run(
        [str(OFFICE_WORKER_BIN), "session"],
        cwd=PROJECT_ROOT,
        input=(
            json.dumps(
                {
                    "contract": "ai_daily_worker",
                    "protocol_version": 2,
                    "request_id": request["request_id"],
                    "operation": "office_parse",
                    "payload": request,
                },
                ensure_ascii=False,
            )
            + "\n"
        ).encode("utf-8"),
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr.decode(
        "utf-8",
        errors="replace",
    )
    frames = completed.stdout.splitlines()
    assert json.loads(frames[0])["frame"] == "hello"
    response = WorkerParseResponse.model_validate(json.loads(frames[1])["result"])
    assert response.status == "ok"
    assert response.parser_backend == backend
    assert response.worker_lane == "rust_office_process_v2"
    assert "strict content" in response.content


def test_office_worker_corrupt_zip_is_deterministic_error(tmp_path: Path) -> None:
    _require_office_worker()
    sample = tmp_path / "corrupt.xlsx"
    sample.write_bytes(b"not a zip")
    request = _office_parse_request(
        sample,
        ".xlsx",
        "rust_xlsx_bounded_v2",
    )

    completed = subprocess.run(
        [str(OFFICE_WORKER_BIN), "session"],
        cwd=PROJECT_ROOT,
        input=(json.dumps({
            "contract": "ai_daily_worker",
            "protocol_version": 2,
            "request_id": request["request_id"],
            "operation": "office_parse",
            "payload": request,
        }) + "\n").encode("utf-8"),
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0
    frames = completed.stdout.splitlines()
    response = WorkerParseResponse.model_validate(json.loads(frames[1])["result"])
    assert response.status == "error"
    assert response.error is not None
    assert response.error.error_code == "PARSER_FAILED"
    assert response.error.retryable is False


def test_office_worker_invalid_v2_request_exits_with_protocol_error() -> None:
    _require_office_worker()

    completed = subprocess.run(
        [str(OFFICE_WORKER_BIN), "session"],
        cwd=PROJECT_ROOT,
        input=b"not-json",
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 2
    assert json.loads(completed.stdout.splitlines()[0])["frame"] == "hello"


def _office_parse_request(path: Path, file_type: str, backend: str) -> dict[str, object]:
    stat = path.stat()
    return {
        "contract": "ai_daily_worker",
        "protocol_version": 1,
        "request_id": "63333333-6333-4333-8333-633333333333",
        "file_path": str(path.resolve()),
        "file_type": file_type,
        "backend": backend,
        "remaining_timeout_ms": 30_000,
        "max_file_size_bytes": 1_000_000,
        "parser_limits": {
            "kind": "office",
            "excel_max_sheets": 2,
            "excel_max_rows": 10,
            "excel_max_columns": 12,
            "docx_max_paragraphs": 80,
            "docx_max_tables": 8,
            "docx_table_max_rows": 20,
            "docx_table_max_cols": 8,
            "pptx_max_slides": 15,
            "pptx_include_notes": True,
            "document_excerpt_max_chars": 4000,
        },
        "expected_source_version": (
            f"mtime_ns={stat.st_mtime_ns}:size={stat.st_size}"
        ),
    }


def _python_document_parse_request(
    path: Path,
    file_type: str,
    backend: str,
) -> dict[str, object]:
    request = _office_parse_request(path, file_type, backend)
    if backend == "python_sharepoint_text_v2":
        request["parser_limits"] = {
            "kind": "sharepoint_text",
            "excerpt_max_chars": 4000,
        }
    return request
