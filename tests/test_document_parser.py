"""测试 Office/PDF bounded document parser。"""

import json
import subprocess
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from src.services.document_parser import (
    OFFICE_PARSER_BACKEND,
    PDF_TEXT_PARSER_BACKEND,
    DocumentParserOptions,
    parse_document_file,
)

RUST_OFFICE_PARSER_BIN = (
    Path(__file__).resolve().parents[1]
    / "rust/office_parser/target/release/ai-daily-office-parser"
)


def _run_rust_office_parser(sample: Path, file_type: str) -> dict:
    request = {
        "file_path": str(sample),
        "file_type": file_type,
        "limits": {"document_excerpt_max_chars": 4000},
        "parser_backend": "rust_office_oxide_v1",
    }
    completed = subprocess.run(
        [str(RUST_OFFICE_PARSER_BIN)],
        input=json.dumps(request, ensure_ascii=False),
        text=True,
        capture_output=True,
        timeout=20,
        check=False,
    )
    assert completed.returncode == 0, completed.stderr
    return json.loads(completed.stdout)


def _write_minimal_pdf(path: Path, content_stream: str = "") -> None:
    """写入最小 PDF fixture，避免为测试引入额外 PDF 生成依赖。"""
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Length "
            + str(len(content_stream.encode("ascii"))).encode("ascii")
            + b" >>\nstream\n"
            + content_stream.encode("ascii")
            + b"\nendstream"
        ),
    ]
    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(body))
        body.extend(f"{index} 0 obj\n".encode("ascii"))
        body.extend(obj)
        body.extend(b"\nendobj\n")
    xref_offset = len(body)
    body.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    body.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(bytes(body))


def test_parse_docx_extracts_paragraphs_and_tables(tmp_path: Path):
    """DOCX preview 应提取段落和表格。"""
    sample = tmp_path / "report.docx"
    doc = Document()
    doc.add_paragraph("项目进展正常")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "事项"
    table.cell(0, 1).text = "状态"
    table.cell(1, 0).text = "日报"
    table.cell(1, 1).text = "完成"
    doc.save(sample)

    context = parse_document_file(
        sample,
        ".docx",
        {"document_excerpt_max_chars": 4000},
        DocumentParserOptions(),
    )

    assert context.error is None
    assert context.parser_backend == OFFICE_PARSER_BACKEND
    assert context.truncated is False
    assert "项目进展正常" in context.content
    assert "| 事项 | 状态 |" in context.content
    assert "| 日报 | 完成 |" in context.content


def test_parse_xlsx_limits_sheets_rows_columns_and_marks_truncated(tmp_path: Path):
    """XLSX preview 应按 sheet/row/column 预算截断。"""
    sample = tmp_path / "workbook.xlsx"
    workbook = Workbook()
    first = workbook.active
    first.title = "SheetA"
    first.append(["A", "B", "C"])
    first.append(["r1a", "r1b", "r1c"])
    first.append(["r2a", "r2b", "r2c"])
    second = workbook.create_sheet("SheetB")
    second.append(["hidden"])
    workbook.save(sample)

    context = parse_document_file(
        sample,
        ".xlsx",
        {
            "excel_max_sheets": 1,
            "excel_max_rows": 2,
            "excel_max_columns": 2,
            "document_excerpt_max_chars": 4000,
        },
        DocumentParserOptions(),
    )

    assert context.error is None
    assert context.parser_backend == OFFICE_PARSER_BACKEND
    assert context.truncated is True
    assert "## Sheet: SheetA" in context.content
    assert "SheetB" not in context.content
    assert "r1a" in context.content
    assert "r2a" not in context.content
    assert "r1c" not in context.content


def test_parse_pptx_extracts_slide_text(tmp_path: Path):
    """PPTX preview 应提取 slide 文本。"""
    sample = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "周报汇总"
    slide.placeholders[1].text = "完成 scanner 优化"
    presentation.save(sample)

    context = parse_document_file(
        sample,
        ".pptx",
        {"pptx_max_slides": 5, "document_excerpt_max_chars": 4000},
        DocumentParserOptions(),
    )

    assert context.error is None
    assert context.parser_backend == OFFICE_PARSER_BACKEND
    assert context.truncated is False
    assert "## Slide 1" in context.content
    assert "周报汇总" in context.content
    assert "完成 scanner 优化" in context.content


@pytest.mark.skipif(
    not RUST_OFFICE_PARSER_BIN.exists(),
    reason="Rust Office parser release binary has not been built",
)
def test_rust_office_parser_extracts_docx_bounded_xlsx_and_pptx(
    tmp_path: Path,
):
    docx_sample = tmp_path / "report.docx"
    doc = Document()
    doc.add_paragraph("Rust DOCX 中文")
    doc.save(docx_sample)

    xlsx_sample = tmp_path / "workbook.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["项目", "状态"])
    sheet.append(["Rust XLSX 中文", "完成"])
    workbook.save(xlsx_sample)

    pptx_sample = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Rust PPTX 中文"
    slide.placeholders[1].text = "完成 parser spike"
    presentation.save(pptx_sample)

    docx_context = _run_rust_office_parser(docx_sample, ".docx")
    xlsx_context = _run_rust_office_parser(xlsx_sample, ".xlsx")
    pptx_context = _run_rust_office_parser(pptx_sample, ".pptx")

    assert docx_context["error"] is None
    assert "Rust DOCX 中文" in docx_context["content"]
    assert xlsx_context["error"] is None
    assert xlsx_context["parser_backend"] == "rust_xlsx_bounded_v1"
    assert "Rust XLSX 中文" in xlsx_context["content"]
    assert pptx_context["error"] is None
    assert "Rust PPTX 中文" in pptx_context["content"]


def test_parse_pdf_extracts_text_layer(tmp_path: Path):
    """文本型 PDF 应提取 text layer。"""
    sample = tmp_path / "text.pdf"
    _write_minimal_pdf(
        sample,
        "BT /F1 24 Tf 100 700 Td (Hello PDF text layer) Tj ET",
    )

    context = parse_document_file(
        sample,
        ".pdf",
        {"pdf_max_pages": 2, "document_excerpt_max_chars": 4000},
        DocumentParserOptions(),
    )

    assert context.error is None
    assert context.parser_backend == PDF_TEXT_PARSER_BACKEND
    assert context.truncated is False
    assert "## Page 1" in context.content
    assert "Hello PDF text layer" in context.content


def test_parse_pdf_without_text_layer_returns_auditable_error(tmp_path: Path):
    """无 text layer 的 PDF 应返回可审计错误，不做 OCR。"""
    sample = tmp_path / "blank.pdf"
    _write_minimal_pdf(sample)

    context = parse_document_file(
        sample,
        ".pdf",
        {"pdf_max_pages": 2, "document_excerpt_max_chars": 4000},
        DocumentParserOptions(),
    )

    assert context.content == ""
    assert context.error is not None
    assert context.error.startswith("PDF_NO_TEXT_LAYER:")
    assert context.parser_backend == PDF_TEXT_PARSER_BACKEND
    assert context.truncated is False
