"""Python document worker 的严格 parse 合同与 legacy adapter 回归。"""

from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from src.workers.contracts import (
    WorkerParsePayload,
    parse_legacy_excel_payload,
    parse_sharepoint_text_payload,
    parse_worker_request,
)
from src.workers.models import (
    OfficeLimits,
    ParseRequest,
    PdfLimits,
    SharePointTextLimits,
    WorkerOperationError,
)


PROJECT_ROOT = Path(__file__).resolve().parents[1]
LEGACY_FIXTURE_DIR = PROJECT_ROOT / "tests" / "fixtures" / "worker_documents"


def _source_version(path: Path) -> str:
    stat = path.stat()
    return f"mtime_ns={stat.st_mtime_ns}:size={stat.st_size}"


def _office_limits() -> OfficeLimits:
    return OfficeLimits(
        kind="office",
        excel_max_sheets=2,
        excel_max_rows=10,
        excel_max_columns=12,
        docx_max_paragraphs=80,
        docx_max_tables=8,
        docx_table_max_rows=20,
        docx_table_max_cols=8,
        pptx_max_slides=15,
        pptx_include_notes=True,
        document_excerpt_max_chars=4000,
    )


def _request(
    path: Path,
    file_type: str,
    backend: str,
    parser_limits: OfficeLimits | PdfLimits | SharePointTextLimits,
    *,
    max_file_size_bytes: int = 1_000_000,
    expected_source_version: str | None = None,
) -> ParseRequest:
    return ParseRequest(
        file_path=str(path.resolve()),
        file_type=file_type,
        backend=backend,
        remaining_timeout_ms=30_000,
        max_file_size_bytes=max_file_size_bytes,
        parser_limits=parser_limits,
        expected_source_version=(
            expected_source_version
            if expected_source_version is not None
            else _source_version(path)
        ),
    )


@pytest.mark.parametrize("file_type", [".docx", ".xlsx", ".pptx"])
def test_python_office_worker_parses_modern_office_fixtures(
    tmp_path: Path,
    file_type: str,
) -> None:
    sample = tmp_path / f"现代 文档{file_type}"
    if file_type == ".docx":
        document = Document()
        document.add_paragraph("DOCX worker content")
        document.save(sample)
    elif file_type == ".xlsx":
        workbook = Workbook()
        workbook.active.append(["XLSX worker content"])
        workbook.save(sample)
        workbook.close()
    else:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "PPTX worker content"
        presentation.save(sample)

    response = parse_worker_request(
        _request(sample, file_type, "python_office_v2", _office_limits())
    )

    assert response.parser_backend == "python_office_v2"
    assert response.worker_lane == "python_document_process_v2"
    assert "worker content" in response.content
    assert response.observed_source_version == _source_version(sample)


def test_pdf_worker_parses_a_bounded_text_layer(tmp_path: Path) -> None:
    sample = tmp_path / "report.pdf"
    _write_minimal_pdf(sample, "BT /F1 12 Tf 72 720 Td (PDF worker text) Tj ET")
    limits = PdfLimits(kind="pdf", max_pages=2, excerpt_max_chars=4000)

    response = parse_worker_request(
        _request(sample, ".pdf", "python_pdf_text_v2", limits)
    )

    assert response.parser_backend == "python_pdf_text_v2"
    assert "PDF worker text" in response.content


def test_worker_size_guard_runs_before_document_parser(tmp_path: Path) -> None:
    sample = tmp_path / "too-large.pdf"
    sample.write_bytes(b"0123456789")
    limits = PdfLimits(kind="pdf", max_pages=2, excerpt_max_chars=4000)

    def forbidden_parser(*args: object, **kwargs: object) -> object:
        raise AssertionError("parser must not run after the size guard")

    with pytest.raises(WorkerOperationError) as captured:
        parse_worker_request(
            _request(
                sample,
                ".pdf",
                "python_pdf_text_v2",
                limits,
                max_file_size_bytes=9,
            ),
            document_parser=forbidden_parser,
        )

    assert captured.value.error_code == "FILE_TOO_LARGE"


def test_worker_rejects_changed_source_before_parser(tmp_path: Path) -> None:
    sample = tmp_path / "changed.pdf"
    sample.write_bytes(b"fixture")
    limits = PdfLimits(kind="pdf", max_pages=2, excerpt_max_chars=4000)

    def forbidden_parser(*args: object, **kwargs: object) -> object:
        raise AssertionError("parser must not run for a stale source version")

    with pytest.raises(WorkerOperationError) as captured:
        parse_worker_request(
            _request(
                sample,
                ".pdf",
                "python_pdf_text_v2",
                limits,
                expected_source_version="mtime_ns=1:size=7",
            ),
            document_parser=forbidden_parser,
        )

    assert captured.value.error_code == "SOURCE_VERSION_CHANGED"
    assert captured.value.retryable is True


def test_worker_diagnostic_does_not_expose_parser_exception_text(
    tmp_path: Path,
) -> None:
    sample = tmp_path / "secret-error.pdf"
    sample.write_bytes(b"fixture")
    limits = PdfLimits(kind="pdf", max_pages=2, excerpt_max_chars=4000)

    def secret_parser(*args: object, **kwargs: object) -> object:
        raise RuntimeError("SECRET_CELL_VALUE must not cross the worker contract")

    with pytest.raises(WorkerOperationError) as captured:
        parse_worker_request(
            _request(sample, ".pdf", "python_pdf_text_v2", limits),
            document_parser=secret_parser,
        )

    assert captured.value.message == "document parser reported an error"
    assert "SECRET_CELL_VALUE" not in str(captured.value)


@pytest.mark.parametrize(
    ("file_name", "file_type", "backend", "expected_text"),
    [
        (
            "legacy_sample.xls",
            ".xls",
            "python_office_v2",
            "Legacy XLS worker content",
        ),
        (
            "legacy_sample.doc",
            ".doc",
            "python_sharepoint_text_v2",
            "Legacy DOC worker content",
        ),
        (
            "legacy_sample.ppt",
            ".ppt",
            "python_sharepoint_text_v2",
            "Legacy PPT worker content",
        ),
    ],
)
def test_strict_python_worker_parses_real_legacy_office_fixtures(
    file_name: str,
    file_type: str,
    backend: str,
    expected_text: str,
) -> None:
    sample = LEGACY_FIXTURE_DIR / file_name
    limits: OfficeLimits | SharePointTextLimits
    if file_type == ".xls":
        limits = _office_limits()
    else:
        limits = SharePointTextLimits(
            kind="sharepoint_text",
            excerpt_max_chars=4000,
        )

    response = parse_worker_request(_request(sample, file_type, backend, limits))

    assert response.parser_backend == backend
    assert response.worker_lane == "python_document_process_v2"
    assert expected_text in response.content
    if file_type == ".xls":
        assert "| Name | Value |" in response.content


def test_legacy_xls_adapter_keeps_table_preview_capability(tmp_path: Path) -> None:
    sample = tmp_path / "legacy.xls"
    sample.write_bytes(b"synthetic legacy fixture")

    payload = parse_legacy_excel_payload(
        sample,
        ".xls",
        {"excel_max_rows": 10, "document_excerpt_max_chars": 20},
        excel_reader=lambda _path, _sheets, rows, _columns: (
            f"## Sheet\n\nlegacy xls rows={rows}",
            False,
        ),
    )

    assert isinstance(payload, WorkerParsePayload)
    assert payload.error is None
    assert payload.parser_backend == "python_office_v2"
    assert payload.content == "## Sheet\n\nlegacy xls"
    assert payload.truncated is True


def test_strict_xls_worker_enforces_sheet_row_and_column_budgets() -> None:
    sample = LEGACY_FIXTURE_DIR / "legacy_sample.xls"
    limits = _office_limits().model_copy(
        update={
            "excel_max_sheets": 1,
            "excel_max_rows": 1,
            "excel_max_columns": 2,
        }
    )

    response = parse_worker_request(
        _request(sample, ".xls", "python_office_v2", limits)
    )

    assert response.truncated is True
    assert "Legacy XLS worker content" in response.content
    assert "Column beyond budget" not in response.content
    assert "Row beyond budget" not in response.content
    assert "Second sheet beyond budget" not in response.content


@pytest.mark.parametrize("file_type", [".doc", ".ppt"])
def test_legacy_sharepoint_adapter_keeps_text_capability(
    tmp_path: Path,
    file_type: str,
) -> None:
    sample = tmp_path / f"legacy{file_type}"
    sample.write_bytes(b"synthetic legacy fixture")

    class FakeResult:
        def get_full_text(self) -> str:
            return "legacy sharepoint text"

    fake_module = SimpleNamespace(read_file=lambda _path: iter([FakeResult()]))
    payload = parse_sharepoint_text_payload(
        sample,
        file_type,
        {"excerpt_max_chars": 12},
        import_module=lambda _name: fake_module,
    )

    assert isinstance(payload, WorkerParsePayload)
    assert payload.error is None
    assert payload.parser_backend == "python_sharepoint_text_v2"
    assert payload.content == "legacy share"
    assert payload.truncated is True


def _write_minimal_pdf(path: Path, content_stream: str) -> None:
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Length "
            + str(len(content_stream.encode("ascii"))).encode("ascii")
            + b" >>\nstream\n"
            + content_stream.encode("ascii")
            + b"\nendstream"
        ),
    ]
    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(body))
        body.extend(f"{index} 0 obj\n".encode("ascii"))
        body.extend(obj)
        body.extend(b"\nendobj\n")
    xref_offset = len(body)
    body.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    body.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(bytes(body))
