"""Generate sanitized fixtures and evaluate Task 10 scanner cutover gates."""

from __future__ import annotations

import argparse
from contextlib import contextmanager
from dataclasses import dataclass, replace
from datetime import date, datetime
import gc
import json
import logging
import math
import os
from pathlib import Path
import shutil
from statistics import median
import subprocess
import sys
import tempfile
import time
from time import perf_counter
from typing import Any, Iterator

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from docx import Document  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from pptx import Presentation  # noqa: E402

import scripts.compare_context_engines as comparison  # noqa: E402
from src.core.config import config as app_config  # noqa: E402
from src.services.context_scheduler import ContextScheduleRequest  # noqa: E402


_ALLOWED_EXTENSIONS = [
    ".csv",
    ".docx",
    ".json",
    ".log",
    ".md",
    ".pdf",
    ".pptx",
    ".txt",
    ".xlsx",
]
_TEXT_AND_PDF_EXTENSIONS = {
    ".csv",
    ".json",
    ".log",
    ".md",
    ".pdf",
    ".txt",
}
_DEFAULT_PROFILE_VERSION = "task10-synthetic-v1"
_CONTEXT_GOLDEN_PATH = (
    PROJECT_ROOT
    / "tests"
    / "fixtures"
    / "scanner_cutover"
    / "task10_expected_context_hashes.json"
)
_REAL_WORK_DIR_ENV = "AI_DAILY_REAL_WORK_DIR"
_REAL_SAMPLE_FILE_ENV = "AI_DAILY_REAL_SAMPLE_FILE"


@dataclass(frozen=True, slots=True)
class GateRuntimeConfig:
    """Scanner-only config clone; never mutates the loaded local settings."""

    scanner_config: dict[str, Any]
    contract_profile: dict[str, Any]
    rust_scanner_bin: str
    rust_process_timeout_seconds: float

    def scanner_contract_profile(self) -> dict[str, Any]:
        return dict(self.contract_profile)

    def with_profile_version(self, value: str) -> "GateRuntimeConfig":
        scanner_config = dict(self.scanner_config)
        scanner_config["parser_profile_version"] = value
        contract_profile = dict(self.contract_profile)
        contract_profile["parser_profile_version"] = value
        return replace(
            self,
            scanner_config=scanner_config,
            contract_profile=contract_profile,
        )


def build_gate_config(base_config: Any) -> GateRuntimeConfig:
    """Clone scanner-only leaves and apply deterministic synthetic limits."""
    shared_overrides: dict[str, Any] = {
        "allowed_extensions": list(_ALLOWED_EXTENSIONS),
        "ignored_patterns": ["slow worker.pdf"],
        "excluded_dirs": ["faults"],
        "max_workers": 2,
        "max_file_size_mb": 1,
        "total_max_chars": 120_000,
        "parser_profile_version": _DEFAULT_PROFILE_VERSION,
        "office_parser_backend": "rust_office_oxide_v1",
        "pdf_parser_backend": "pdf_text_v1",
        "office_fallback_after_timeout": False,
        "office_legacy_extensions_enabled": False,
    }
    scanner_config = dict(base_config.scanner_config)
    scanner_config.update(shared_overrides)
    contract_profile = dict(base_config.scanner_contract_profile())
    contract_profile.update(shared_overrides)
    contract_profile["schema_version"] = "scanner_profile_v1"
    return GateRuntimeConfig(
        scanner_config=scanner_config,
        contract_profile=contract_profile,
        rust_scanner_bin=str(base_config.rust_scanner_bin),
        rust_process_timeout_seconds=float(
            base_config.rust_process_timeout_seconds
        ),
    )


def build_real_gate_config(base_config: Any) -> GateRuntimeConfig:
    """Clone the effective scanner leaves without changing local config."""
    return GateRuntimeConfig(
        scanner_config=dict(base_config.scanner_config),
        contract_profile=dict(base_config.scanner_contract_profile()),
        rust_scanner_bin=str(base_config.rust_scanner_bin),
        rust_process_timeout_seconds=float(
            base_config.rust_process_timeout_seconds
        ),
    )


def create_synthetic_corpus(
    work_dir: Path,
    *,
    scan_date: date | None = None,
) -> dict[str, Any]:
    """Create the complete sanitized corpus without reading project data."""
    work_dir = work_dir.resolve(strict=True)
    if not work_dir.is_dir() or any(work_dir.iterdir()):
        raise ValueError("synthetic work directory must exist and be empty")

    text_files = {
        "plain.txt": "Synthetic plain text evidence.\n",
        "notes with space.md": "# Synthetic notes\n\nDeterministic evidence.\n",
        "table.csv": "name,value\nalpha,1\nbeta,2\n",
        "data.json": '{"alpha":1,"items":["a","b"]}\n',
        "events.log": "".join(
            f"synthetic-event-{index:03d}\n" for index in range(200)
        ),
    }
    for relative_path, content in text_files.items():
        (work_dir / relative_path).write_text(content, encoding="utf-8")

    unicode_dir = work_dir / "\u4e2d\u6587 \u8def\u5f84"
    unicode_dir.mkdir()
    xlsx_path = unicode_dir / "synthetic workbook.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet.append(["item", "value"])
    sheet.append(["synthetic", 42])
    workbook.save(xlsx_path)
    workbook.close()

    docx_path = work_dir / "synthetic document.docx"
    document = Document()
    document.add_paragraph("Synthetic DOCX evidence")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "item"
    table.cell(0, 1).text = "value"
    table.cell(1, 0).text = "synthetic"
    table.cell(1, 1).text = "42"
    document.save(docx_path)

    pptx_path = work_dir / "synthetic slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Synthetic presentation"
    slide.placeholders[1].text = "Deterministic PPTX evidence"
    presentation.save(pptx_path)

    pdf_path = work_dir / "synthetic evidence.pdf"
    _write_minimal_pdf(
        pdf_path,
        "BT /F1 12 Tf 72 720 Td (Synthetic PDF evidence) Tj ET",
    )

    corrupt_path = work_dir / "corrupt workbook.xlsx"
    corrupt_path.write_bytes(b"synthetic invalid Office ZIP")
    oversized_path = work_dir / "oversized.txt"
    with oversized_path.open("wb") as handle:
        handle.truncate(1_048_577)

    fault_dir = work_dir / "faults"
    fault_dir.mkdir()
    slow_path = fault_dir / "slow worker.pdf"
    _write_minimal_pdf(
        slow_path,
        "BT /F1 12 Tf 72 720 Td (Synthetic slow worker fixture) Tj ET",
    )

    if scan_date is not None:
        fixture_timestamp = datetime.combine(
            scan_date,
            datetime.min.time(),
        ).replace(hour=12).timestamp()
        for path in work_dir.rglob("*"):
            if path.is_file():
                os.utime(path, (fixture_timestamp, fixture_timestamp))

    scanned_files = sorted(
        path.relative_to(work_dir).as_posix()
        for path in work_dir.rglob("*")
        if path.is_file() and fault_dir not in path.parents
    )
    return {
        "scanned_files": scanned_files,
        "corrupt_office_count": 1,
        "oversized_file_count": 1,
        "slow_worker_fixture_count": 1,
    }


def run_parity_gate(
    *,
    gate_root: Path,
    gate_config: GateRuntimeConfig,
    scan_date: date,
) -> dict[str, Any]:
    """Run two cold complete-context comparisons and evaluate parity."""
    work_dir = _new_directory(gate_root / "work")
    manifest = create_synthetic_corpus(work_dir, scan_date=scan_date)
    payloads: list[dict[str, Any]] = []
    with _using_gate_config(gate_config):
        for index in range(2):
            shadow_root = _new_directory(gate_root / f"shadow-{index + 1}")
            payloads.append(
                comparison.compare_context_engines(
                    work_dir=work_dir,
                    start_date=scan_date,
                    end_date=scan_date,
                    report_mode="daily",
                    redact_content=True,
                    ephemeral_db_root=shadow_root,
                    output=shadow_root / "comparison.json",
                )
            )

    first, second = payloads
    legacy = first["engines"]["python_legacy"]
    rust = first["engines"]["rust_v2"]
    legacy_summary = legacy["summary"]
    rust_summary = rust["summary"]
    count_fields = ("success_count", "error_file_count", "timeout_count")
    counts_equal = all(
        legacy_summary[field] == rust_summary[field] for field in count_fields
    )
    text_pdf_differences = [
        item
        for item in first["content_hash_differences"]
        if Path(item["relative_path"]).suffix.lower()
        in _TEXT_AND_PDF_EXTENSIONS
    ]
    same_backend_differences = [
        item
        for item in first["content_hash_differences"]
        if item["python_legacy_backend"] == item["rust_v2_backend"]
    ]
    deterministic = {
        engine: (
            first["engines"][engine]["context_sha256"]
            == second["engines"][engine]["context_sha256"]
        )
        for engine in ("python_legacy", "rust_v2")
    }
    normalized_context_hashes = {
        engine: str(
            first["engines"][engine]["normalized_context_sha256"]
        )
        for engine in ("python_legacy", "rust_v2")
    }
    context_golden = _load_context_golden()
    golden_matches = (
        normalized_context_hashes
        == context_golden["normalized_context_sha256"]
        and len(manifest["scanned_files"])
        == context_golden["source_file_count"]
    )
    total_max_chars = int(gate_config.contract_profile["total_max_chars"])
    within_budget = all(
        first["engines"][engine]["summary"]["output_chars"]
        <= total_max_chars
        for engine in ("python_legacy", "rust_v2")
    )
    backend_lane_evidence = all(
        item.get("parser_backend") and item.get("worker_lane")
        for engine in ("python_legacy", "rust_v2")
        for item in first["engines"][engine]["files"]
    )
    expected_inventory_count = len(manifest["scanned_files"])
    inventory_count_matches = all(
        first["engines"][engine]["summary"]["source_file_count"]
        == expected_inventory_count
        for engine in ("python_legacy", "rust_v2")
    )
    violations: list[str] = []
    _require(violations, first["inventory_difference_count"] == 0, "inventory")
    _require(violations, inventory_count_matches, "inventory_count")
    _require(violations, counts_equal, "status_counts")
    _require(violations, not text_pdf_differences, "text_pdf_hashes")
    _require(violations, not same_backend_differences, "same_backend_hashes")
    _require(
        violations,
        first["decision_difference_count"] == 0,
        "decisions",
    )
    _require(violations, all(deterministic.values()), "determinism")
    _require(violations, golden_matches, "intentional_context_difference_golden")
    _require(violations, within_budget, "context_budget")
    _require(violations, backend_lane_evidence, "backend_lane_evidence")
    _require(
        violations,
        legacy["status"] != "error" and rust["status"] != "error",
        "engine_status",
    )
    _require(violations, first["fallback_count"] == 0, "fallback")
    return {
        "passed": not violations,
        "violations": violations,
        "inventory_difference_count": first["inventory_difference_count"],
        "content_hash_difference_count": first[
            "content_hash_difference_count"
        ],
        "text_pdf_hash_difference_count": len(text_pdf_differences),
        "same_backend_hash_difference_count": len(same_backend_differences),
        "decision_difference_count": first["decision_difference_count"],
        "fallback_count": first["fallback_count"],
        "counts_equal": counts_equal,
        "within_budget": within_budget,
        "backend_lane_evidence": backend_lane_evidence,
        "engine_summaries": {
            "python_legacy": legacy_summary,
            "rust_v2": rust_summary,
        },
        "inventory_differences": first["inventory_differences"],
        "decision_differences": first["decision_differences"],
        "cross_engine_final_context_equal": first[
            "final_context_hash_equal"
        ],
        "normalized_context_sha256": normalized_context_hashes,
        "intentional_context_difference_golden_matches": golden_matches,
        "deterministic": deterministic,
        "source_file_count": expected_inventory_count,
    }


def run_cache_gate(
    *,
    gate_root: Path,
    gate_config: GateRuntimeConfig,
    scan_date: date,
) -> dict[str, Any]:
    """Prove cold, warm, one-file, and semantic-profile cache behavior."""
    work_dir = _new_directory(gate_root / "work")
    create_synthetic_corpus(work_dir, scan_date=scan_date)
    db_root = _new_directory(gate_root / "databases")
    legacy_db_root = _new_directory(db_root / "python_legacy")
    rust_db_root = _new_directory(db_root / "rust_v2")
    db_paths = {
        "python_legacy": legacy_db_root / "scan.sqlite3",
        "rust_v2": rust_db_root / "scan_index_v2.sqlite3",
    }
    request = _request(scan_date)
    observations: dict[str, dict[str, dict[str, Any]]] = {
        "python_legacy": {},
        "rust_v2": {},
    }
    with _using_gate_config(gate_config):
        for engine in observations:
            observations[engine]["cold"] = _run_engine(
                engine,
                request,
                work_dir,
                db_paths[engine],
            )
            observations[engine]["warm"] = _run_engine(
                engine,
                request,
                work_dir,
                db_paths[engine],
            )

        mutation_file = work_dir / "plain.txt"
        mutation_file.write_text(
            mutation_file.read_text(encoding="utf-8")
            + "Synthetic one-file mutation.\n",
            encoding="utf-8",
        )
        for engine in observations:
            observations[engine]["single_file"] = _run_engine(
                engine,
                request,
                work_dir,
                db_paths[engine],
            )

    changed_config = gate_config.with_profile_version(
        "task10-synthetic-profile-change-v2"
    )
    with _using_gate_config(changed_config):
        for engine in observations:
            observations[engine]["profile_change"] = _run_engine(
                engine,
                request,
                work_dir,
                db_paths[engine],
            )

    successful_counts = {
        engine: len(_successful_files(runs["cold"]))
        for engine, runs in observations.items()
    }
    successful_path_sets = {
        engine: _successful_path_keys(runs["cold"])
        for engine, runs in observations.items()
    }
    successful_file_count = successful_counts["python_legacy"]
    mutation_key = "plain.txt"
    violations: list[str] = []
    _require(
        violations,
        len(set(successful_counts.values())) == 1,
        "successful_file_count",
    )
    _require(
        violations,
        len({frozenset(paths) for paths in successful_path_sets.values()}) == 1,
        "successful_file_identity",
    )
    engine_evidence: dict[str, dict[str, Any]] = {}
    for engine, runs in observations.items():
        successful_paths = successful_path_sets[engine]
        cold_miss_paths = _cache_path_keys(runs["cold"], "miss")
        warm_fresh_paths = _cache_path_keys(runs["warm"], "fresh")
        single_miss_paths = _cache_path_keys(runs["single_file"], "miss")
        single_fresh_paths = _cache_path_keys(runs["single_file"], "fresh")
        profile_miss_paths = _cache_path_keys(runs["profile_change"], "miss")
        single_file_identity_matched = (
            single_miss_paths == {mutation_key}
            and single_fresh_paths == successful_paths - {mutation_key}
        )
        evidence = {
            "cold_successful_misses": len(cold_miss_paths),
            "warm_successful_fresh": len(warm_fresh_paths),
            "single_file_successful_misses": len(single_miss_paths),
            "single_file_successful_fresh": len(single_fresh_paths),
            "single_file_identity_matched": single_file_identity_matched,
            "profile_change_successful_misses": len(profile_miss_paths),
        }
        engine_evidence[engine] = evidence
        _require(
            violations,
            evidence["cold_successful_misses"] == successful_file_count,
            f"{engine}_cold",
        )
        _require(
            violations,
            evidence["warm_successful_fresh"] == successful_file_count,
            f"{engine}_warm",
        )
        _require(
            violations,
            evidence["single_file_successful_misses"] == 1,
            f"{engine}_single_file",
        )
        _require(
            violations,
            single_file_identity_matched,
            f"{engine}_single_file_identity",
        )
        _require(
            violations,
            evidence["profile_change_successful_misses"]
            == successful_file_count,
            f"{engine}_profile_change",
        )
        _require(
            violations,
            cold_miss_paths == successful_paths
            and warm_fresh_paths == successful_paths
            and profile_miss_paths == successful_paths,
            f"{engine}_cache_file_identity",
        )
        _require(
            violations,
            runs["cold"]["context_sha256"]
            == runs["warm"]["context_sha256"],
            f"{engine}_warm_determinism",
        )
    return {
        "passed": not violations,
        "violations": violations,
        "successful_file_count": successful_file_count,
        "engines": engine_evidence,
    }


def run_performance_gate(
    *,
    gate_root: Path,
    gate_config: GateRuntimeConfig,
    scan_date: date,
    samples: int,
) -> dict[str, Any]:
    """Capture at least five cold/warm complete-context samples per engine."""
    if samples < 5:
        raise ValueError("performance gate requires at least five samples")
    gate_started = perf_counter()
    processes_before = _scanner_process_snapshot()
    try:
        collected = _collect_performance_samples(
            gate_root=gate_root,
            gate_config=gate_config,
            scan_date=scan_date,
            samples=samples,
        )
        performance = summarize_performance(
            legacy_cold=collected["durations"]["python_legacy"]["cold"],
            legacy_warm=collected["durations"]["python_legacy"]["warm"],
            rust_cold=collected["durations"]["rust_v2"]["cold"],
            rust_warm=collected["durations"]["rust_v2"]["warm"],
        )
        violations = collected["violations"]
        _require(violations, performance["passed"], "performance_thresholds")
    finally:
        processes_after, processes_settled = _wait_for_scanner_processes(
            processes_before,
        )
    process_evidence = _process_delta_evidence(
        processes_before,
        processes_after,
    )
    _require(violations, processes_settled, "orphan_process")
    elapsed_seconds = perf_counter() - gate_started
    maximum_seconds = max(300.0, samples * 4 * 60.0)
    completed_without_freeze = elapsed_seconds <= maximum_seconds
    _require(violations, completed_without_freeze, "performance_freeze")
    return {
        "passed": not violations,
        "violations": violations,
        "samples_per_engine": samples,
        "performance": performance,
        "parser_backend_counts": collected["parser_backend_counts"],
        "worker_lane_counts": collected["worker_lane_counts"],
        "completed_without_freeze": completed_without_freeze,
        "elapsed_seconds": round(elapsed_seconds, 3),
        "maximum_seconds": maximum_seconds,
        "processes": process_evidence,
    }


def _collect_performance_samples(
    *,
    gate_root: Path,
    gate_config: GateRuntimeConfig,
    scan_date: date,
    samples: int,
) -> dict[str, Any]:
    work_dir = _new_directory(gate_root / "work")
    create_synthetic_corpus(work_dir, scan_date=scan_date)
    db_root = _new_directory(gate_root / "databases")
    request = _request(scan_date)
    durations: dict[str, dict[str, list[float]]] = {
        engine: {"cold": [], "warm": []}
        for engine in ("python_legacy", "rust_v2")
    }
    backend_counts: dict[str, dict[str, int]] = {
        "python_legacy": {},
        "rust_v2": {},
    }
    lane_counts: dict[str, dict[str, int]] = {
        "python_legacy": {},
        "rust_v2": {},
    }
    violations: list[str] = []
    db_paths: list[Path] = []
    with _using_gate_config(gate_config):
        for index in range(samples):
            order = (
                ("python_legacy", "rust_v2")
                if index % 2 == 0
                else ("rust_v2", "python_legacy")
            )
            for engine in order:
                sample_root = _new_directory(
                    db_root / engine / f"sample-{index + 1}"
                )
                db_filename = (
                    "scan_index_v2.sqlite3"
                    if engine == "rust_v2"
                    else "scan.sqlite3"
                )
                db_path = sample_root / db_filename
                db_paths.append(db_path.resolve(strict=False))
                cold = _run_engine(engine, request, work_dir, db_path)
                warm = _run_engine(engine, request, work_dir, db_path)
                durations[engine]["cold"].append(
                    float(cold["build_wall_duration_ms"])
                )
                durations[engine]["warm"].append(
                    float(warm["build_wall_duration_ms"])
                )
                successful_count = len(_successful_files(cold))
                _require(
                    violations,
                    _cache_count(cold, "miss") == successful_count,
                    f"{engine}_sample_{index + 1}_cold_cache",
                )
                _require(
                    violations,
                    _cache_count(warm, "fresh") == successful_count,
                    f"{engine}_sample_{index + 1}_warm_cache",
                )
                _require(
                    violations,
                    cold["status"] != "error" and warm["status"] != "error",
                    f"{engine}_sample_{index + 1}_status",
                )
                for item in cold["files"]:
                    _increment(backend_counts[engine], item["parser_backend"])
                    _increment(lane_counts[engine], item["worker_lane"])

    _require(
        violations,
        len(db_paths) == len(set(db_paths)) == samples * 2,
        "distinct_cold_databases",
    )
    _require(
        violations,
        all(backend_counts[engine] for engine in backend_counts),
        "parser_backend_evidence",
    )
    _require(
        violations,
        all(lane_counts[engine] for engine in lane_counts),
        "worker_lane_evidence",
    )
    return {
        "violations": violations,
        "durations": durations,
        "parser_backend_counts": backend_counts,
        "worker_lane_counts": lane_counts,
    }


def run_fault_gate(*, gate_root: Path) -> dict[str, Any]:
    """Execute the frozen fault suites without serializing captured output."""
    fault_started = perf_counter()
    processes_before = _scanner_process_snapshot()
    python_basetemp = _new_directory(gate_root / "python-basetemp")
    checks = [
        {
            "name": "python_transport_and_no_top_level_fallback",
            "cwd": PROJECT_ROOT,
            "timeout_seconds": 240.0,
            "coverage": [
                "scanner_executable_missing",
                "contract_mismatch",
                "malformed_request_json",
                "malformed_response_json",
                "invalid_utf8_response",
                "whole_process_timeout",
                "python_worker_handshake_invalid",
                "no_top_level_python_fallback",
            ],
            "command": [
                sys.executable,
                "-m",
                "pytest",
                (
                    "tests/test_rust_context_client.py::"
                    "test_json_process_client_maps_timeout_without_trusting_output"
                ),
                (
                    "tests/test_rust_context_client.py::"
                    "test_json_process_client_maps_missing_executable_to_start_failure"
                ),
                (
                    "tests/test_rust_context_client.py::"
                    "test_json_process_client_rejects_nonzero_exit_with_invalid_json"
                ),
                (
                    "tests/test_rust_context_client.py::"
                    "test_json_process_client_rejects_invalid_utf8_stdout"
                ),
                (
                    "tests/test_rust_context_client.py::"
                    "test_json_process_client_rejects_contract_version_mismatch"
                ),
                (
                    "tests/test_rust_context_client.py::"
                    "test_rust_scanner_invalid_json_returns_transport_error_exit_two"
                ),
                (
                    "tests/test_rust_context_client.py::"
                    "test_rust_scanner_doctor_rejects_invalid_python_worker_handshake"
                ),
                (
                    "tests/test_context_builder.py::"
                    "test_error_envelope_is_returned_without_any_legacy_fallback_call"
                ),
                "-q",
                "--basetemp",
                str(python_basetemp),
            ],
            "expected_passed_count": 15,
        },
        _cargo_fault_spec(
            "missing_workers",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="missing_office_and_python_workers_are_environment_unavailable",
            coverage=[
                "office_worker_missing",
                "python_worker_missing",
            ],
        ),
        _cargo_fault_spec(
            "worker_malformed_json",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="invalid_json_is_a_contract_failure",
            coverage=["worker_malformed_json"],
        ),
        _cargo_fault_spec(
            "worker_contract_path_backend",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="wrong_path_or_backend_is_a_contract_failure",
            coverage=["worker_contract_mismatch"],
        ),
        _cargo_fault_spec(
            "worker_contract_build",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="changed_build_after_preflight_is_a_contract_failure",
            coverage=["worker_build_changed"],
        ),
        _cargo_fault_spec(
            "worker_timeout",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="sleep_past_deadline_is_deterministic_timeout",
            coverage=["worker_timeout"],
        ),
        _cargo_fault_spec(
            "worker_crash",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="worker_crash_is_recoverable_parser_failure",
            coverage=["worker_crash"],
        ),
        _cargo_fault_spec(
            "worker_timeout_no_fallback",
            package="ai-daily-scanner-core",
            target="worker_routing",
            test_name="timeout_does_not_fallback_by_default",
            coverage=["timeout_no_fallback"],
        ),
        _cargo_fault_spec(
            "worker_start_failure",
            package="ai-daily-scanner-core",
            target="worker_process",
            test_name="missing_worker_is_an_explicit_start_failure",
            coverage=["worker_start_failure"],
        ),
        _cargo_fault_spec(
            "worker_timeout_process_tree",
            package="ai-daily-scanner-core",
            target="worker_process",
            test_name="timeout_terminates_worker_grandchild",
            coverage=["worker_timeout_process_tree"],
        ),
        _cargo_fault_spec(
            "python_watchdog_process_tree",
            package="ai-daily-scanner-core",
            target="worker_process",
            test_name="python_outer_watchdog_closes_job_and_kills_grandchild",
            coverage=["python_watchdog_process_tree"],
        ),
        _cargo_fault_spec(
            "worker_output_bounds",
            package="ai-daily-scanner-core",
            target="worker_process",
            test_name="output_past_capture_limit_is_rejected_without_pipe_deadlock",
            coverage=["worker_output_bounds"],
        ),
        _cargo_fault_spec(
            "corrupt_office_contract",
            package="ai-daily-office-parser",
            target="worker_contract",
            test_name="corrupt_xlsx_is_a_non_retryable_structured_failure",
            coverage=["corrupt_workbook"],
        ),
        _cargo_fault_spec(
            "unreadable_discovery_entry",
            package="ai-daily-discovery",
            target=None,
            test_name="tests::candidate_metadata_failure_is_a_structured_issue",
            coverage=["unreadable_entry"],
        ),
        _cargo_fault_spec(
            "sqlite_lock",
            package="ai-daily-scanner-core",
            target=None,
            test_name=(
                "store::tests::"
                "database_lock_maps_to_a_structured_retryable_error"
            ),
            coverage=["sqlite_lock"],
        ),
        _cargo_fault_spec(
            "interrupted_transaction",
            package="ai-daily-scanner-core",
            target=None,
            test_name=(
                "store::tests::"
                "failed_final_transaction_leaves_no_cache_inventory_or_false_success"
            ),
            coverage=["interrupted_transaction"],
        ),
    ]
    prerequisites = {
        "windows": sys.platform == "win32",
        "project_venv_python": (
            PROJECT_ROOT / ".venv" / "Scripts" / "python.exe"
        ).is_file(),
    }
    if all(prerequisites.values()):
        results = [_run_fault_check(spec) for spec in checks]
    else:
        results = [
            {
                "name": spec["name"],
                "passed": False,
                "return_code": None,
                "failure_kind": "prerequisite_failed",
                "duration_ms": 0.0,
                "stdout_bytes": 0,
                "stderr_bytes": 0,
                "expected_passed_count": spec["expected_passed_count"],
            }
            for spec in checks
        ]
    processes_after, processes_settled = _wait_for_scanner_processes(
        processes_before,
    )
    process_evidence = _process_delta_evidence(
        processes_before,
        processes_after,
    )
    coverage = {
        coverage_name: result["passed"]
        for spec, result in zip(checks, results, strict=True)
        for coverage_name in spec["coverage"]
    }
    return {
        "passed": all(result["passed"] for result in results)
        and processes_settled,
        "checks": results,
        "prerequisites": prerequisites,
        "coverage": coverage,
        "processes": process_evidence,
        "elapsed_seconds": round(perf_counter() - fault_started, 3),
    }


def run_real_directory_gate(
    *,
    gate_root: Path,
    gate_config: GateRuntimeConfig,
    work_dir: Path,
    sample_file: Path,
) -> dict[str, Any]:
    """Compare the configured local snapshot while retaining aggregate data only."""
    work_path = work_dir.resolve(strict=True)
    sample_path = sample_file.resolve(strict=True)
    if not work_path.is_dir() or not sample_path.is_file():
        raise ValueError("real-directory gate paths are invalid")
    if not sample_path.is_relative_to(work_path):
        raise ValueError("real sample must stay inside the real work directory")
    scan_date = datetime.fromtimestamp(sample_path.stat().st_mtime).date()
    gate_root.resolve(strict=True)
    with _temporary_real_comparison_root() as comparison_root:
        with _using_gate_config(gate_config):
            payload = comparison.compare_context_engines(
                work_dir=work_path,
                start_date=scan_date,
                end_date=scan_date,
                report_mode="daily",
                redact_content=True,
                ephemeral_db_root=comparison_root,
                output=comparison_root / "comparison.json",
            )

    legacy = payload["engines"]["python_legacy"]
    rust = payload["engines"]["rust_v2"]
    rust_xlsx = [
        item
        for item in rust["files"]
        if Path(item["relative_path"]).suffix.lower() == ".xlsx"
    ]
    xlsx_backend_counts: dict[str, int] = {}
    xlsx_lane_counts: dict[str, int] = {}
    for item in rust_xlsx:
        _increment(xlsx_backend_counts, str(item["parser_backend"]))
        _increment(xlsx_lane_counts, str(item["worker_lane"]))

    violations: list[str] = []
    _require(
        violations,
        payload["inventory_difference_count"] == 0,
        "real_inventory",
    )
    _require(
        violations,
        payload["content_hash_difference_count"] == 0,
        "real_content_hashes",
    )
    _require(
        violations,
        payload["decision_difference_count"] == 0,
        "real_decisions",
    )
    _require(violations, payload["fallback_count"] == 0, "real_fallback")
    _require(
        violations,
        legacy["status"] != "error" and rust["status"] != "error",
        "real_engine_status",
    )
    _require(
        violations,
        legacy["summary"]["source_file_count"]
        == rust["summary"]["source_file_count"],
        "real_inventory_count",
    )
    _require(violations, bool(rust_xlsx), "real_xlsx_present")
    _require(
        violations,
        set(xlsx_backend_counts) == {"rust_xlsx_bounded_v1"},
        "real_xlsx_backend",
    )
    _require(
        violations,
        set(xlsx_lane_counts) == {"rust_office_process"},
        "real_xlsx_lane",
    )
    evidence = {
        "passed": not violations,
        "violations": violations,
        "inventory_difference_count": payload[
            "inventory_difference_count"
        ],
        "content_hash_difference_count": payload[
            "content_hash_difference_count"
        ],
        "decision_difference_count": payload["decision_difference_count"],
        "fallback_count": payload["fallback_count"],
        "source_file_count": rust["summary"]["source_file_count"],
        "xlsx_count": len(rust_xlsx),
        "xlsx_parser_backend_counts": xlsx_backend_counts,
        "xlsx_worker_lane_counts": xlsx_lane_counts,
        "engine_statuses": {
            "python_legacy": legacy["status"],
            "rust_v2": rust["status"],
        },
    }
    comparison._assert_metadata_only(evidence)
    return evidence


def summarize_performance(
    *,
    legacy_cold: list[float],
    legacy_warm: list[float],
    rust_cold: list[float],
    rust_warm: list[float],
) -> dict[str, Any]:
    """Apply the frozen median/p95 regression thresholds."""
    samples = [legacy_cold, legacy_warm, rust_cold, rust_warm]
    if any(len(values) < 5 for values in samples):
        raise ValueError("performance summary requires five samples per series")
    stats = {
        "python_legacy": {
            "cold_median_ms": median(legacy_cold),
            "cold_p95_ms": _percentile(legacy_cold, 0.95),
            "warm_median_ms": median(legacy_warm),
        },
        "rust_v2": {
            "cold_median_ms": median(rust_cold),
            "cold_p95_ms": _percentile(rust_cold, 0.95),
            "warm_median_ms": median(rust_warm),
        },
    }
    legacy = stats["python_legacy"]
    rust = stats["rust_v2"]
    criteria = {
        "cold_median_within_10_percent": (
            rust["cold_median_ms"] <= legacy["cold_median_ms"] * 1.10
        ),
        "cold_p95_within_20_percent": (
            rust["cold_p95_ms"] <= legacy["cold_p95_ms"] * 1.20
        ),
        "warm_median_no_regression": (
            rust["warm_median_ms"] <= legacy["warm_median_ms"]
        ),
    }
    rounded_stats = {
        engine: {key: round(float(value), 3) for key, value in values.items()}
        for engine, values in stats.items()
    }
    return {
        "passed": all(criteria.values()),
        "criteria": criteria,
        "stats": rounded_stats,
    }


def _write_minimal_pdf(path: Path, content_stream: str) -> None:
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Length "
            + str(len(content_stream.encode("ascii"))).encode("ascii")
            + b" >>\nstream\n"
            + content_stream.encode("ascii")
            + b"\nendstream"
        ),
    ]
    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, item in enumerate(objects, start=1):
        offsets.append(len(body))
        body.extend(f"{index} 0 obj\n".encode("ascii"))
        body.extend(item)
        body.extend(b"\nendobj\n")
    xref_offset = len(body)
    body.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    body.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(bytes(body))


@contextmanager
def _using_gate_config(gate_config: GateRuntimeConfig) -> Iterator[None]:
    original = comparison.config
    comparison.config = gate_config
    try:
        yield
    finally:
        comparison.config = original


def _request(scan_date: date) -> ContextScheduleRequest:
    return ContextScheduleRequest(
        report_mode="daily",
        source="scan",
        start_date=scan_date,
        end_date=scan_date,
    )


def _run_engine(
    engine: str,
    request: ContextScheduleRequest,
    work_dir: Path,
    db_path: Path,
) -> dict[str, Any]:
    if engine == "python_legacy":
        return comparison.run_legacy_shadow(request, work_dir, db_path)
    if engine == "rust_v2":
        return comparison.run_rust_shadow(request, work_dir, db_path)
    raise ValueError(f"unsupported gate engine: {engine!r}")


def _successful_files(observation: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        item
        for item in observation["files"]
        if item["parse_status"] == "success"
    ]


def _cache_count(observation: dict[str, Any], status: str) -> int:
    return sum(
        item["cache_status"] == status
        for item in _successful_files(observation)
    )


def _successful_path_keys(observation: dict[str, Any]) -> set[str]:
    return {
        comparison._normalize_relative(str(item["relative_path"])).casefold()
        for item in _successful_files(observation)
    }


def _cache_path_keys(
    observation: dict[str, Any],
    status: str,
) -> set[str]:
    return {
        comparison._normalize_relative(str(item["relative_path"])).casefold()
        for item in _successful_files(observation)
        if item["cache_status"] == status
    }


def _load_context_golden() -> dict[str, Any]:
    payload = json.loads(_CONTEXT_GOLDEN_PATH.read_text(encoding="utf-8"))
    expected_keys = {
        "contract",
        "protocol_version",
        "source_file_count",
        "normalized_context_sha256",
    }
    if set(payload) != expected_keys:
        raise ValueError("Task 10 context golden has unexpected fields")
    if (
        payload["contract"] != "ai_daily_task10_context_golden"
        or payload["protocol_version"] != 1
        or not isinstance(payload["source_file_count"], int)
        or set(payload["normalized_context_sha256"])
        != {"python_legacy", "rust_v2"}
    ):
        raise ValueError("Task 10 context golden is invalid")
    for value in payload["normalized_context_sha256"].values():
        if not isinstance(value, str) or len(value) != 64 or any(
            char not in "0123456789abcdef" for char in value
        ):
            raise ValueError("Task 10 context golden hash is invalid")
    return payload


def _cargo_fault_spec(
    name: str,
    *,
    package: str,
    target: str | None,
    test_name: str,
    coverage: list[str],
) -> dict[str, Any]:
    command = ["cargo", "test", "-p", package]
    if target is not None:
        command.extend(["--test", target])
    command.extend([test_name, "--", "--exact"])
    return {
        "name": name,
        "cwd": PROJECT_ROOT / "rust",
        "timeout_seconds": 240.0,
        "coverage": coverage,
        "command": command,
        "expected_passed_count": 1,
    }


def _run_fault_check(spec: dict[str, Any]) -> dict[str, Any]:
    started = perf_counter()
    process: subprocess.Popen[bytes] | None = None
    stdout = b""
    stderr = b""
    failure_kind: str | None = None
    tree_terminated: bool | None = None
    try:
        process = subprocess.Popen(
            [str(value) for value in spec["command"]],
            cwd=str(spec["cwd"]),
            stdin=subprocess.DEVNULL,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            creationflags=(
                getattr(subprocess, "CREATE_NO_WINDOW", 0)
                | getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0)
                if sys.platform == "win32"
                else 0
            ),
        )
        try:
            stdout, stderr = process.communicate(
                timeout=float(spec["timeout_seconds"]),
            )
        except subprocess.TimeoutExpired:
            failure_kind = "timeout"
            tree_terminated = _terminate_process_tree(process)
            stdout, stderr = process.communicate(timeout=15.0)
        if process.returncode != 0 and failure_kind is None:
            failure_kind = "nonzero_exit"
    except OSError:
        failure_kind = "start_failure"
    finally:
        if process is not None and process.poll() is None:
            cleaned = _terminate_process_tree(process)
            tree_terminated = (
                cleaned
                if tree_terminated is None
                else tree_terminated and cleaned
            )
    expected_passed_count = int(spec["expected_passed_count"])
    passed_marker = f"{expected_passed_count} passed".encode("ascii")
    test_count_verified = passed_marker in stdout + stderr
    if (
        process is not None
        and process.returncode == 0
        and failure_kind is None
        and not test_count_verified
    ):
        failure_kind = "test_count_mismatch"
    passed = (
        process is not None
        and process.returncode == 0
        and failure_kind is None
        and test_count_verified
    )
    return {
        "name": spec["name"],
        "passed": passed,
        "return_code": None if process is None else process.returncode,
        "failure_kind": failure_kind,
        "expected_passed_count": expected_passed_count,
        "test_count_verified": test_count_verified,
        "timeout_tree_terminated": tree_terminated,
        "duration_ms": round((perf_counter() - started) * 1000, 3),
        "stdout_bytes": len(stdout),
        "stderr_bytes": len(stderr),
    }


def _terminate_process_tree(process: subprocess.Popen[bytes]) -> bool:
    if process.poll() is not None:
        return True
    if sys.platform == "win32":
        try:
            subprocess.run(
                [
                    "taskkill.exe",
                    "/PID",
                    str(process.pid),
                    "/T",
                    "/F",
                ],
                stdin=subprocess.DEVNULL,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=15.0,
                check=False,
                creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
            )
        except (OSError, subprocess.TimeoutExpired):
            pass
    if process.poll() is None:
        process.kill()
    try:
        process.wait(timeout=15.0)
    except subprocess.TimeoutExpired:
        return False
    return process.poll() is not None


def _scanner_process_snapshot() -> dict[str, list[int]]:
    if sys.platform != "win32":
        raise RuntimeError("Task 10 process evidence requires Windows")
    script = (
        "$all=@(Get-CimInstance Win32_Process -ErrorAction Stop);"
        "$scanner=@($all|Where-Object {$_.Name -eq "
        "'ai-daily-scanner.exe'}|ForEach-Object {[int]$_.ProcessId});"
        "$office=@($all|Where-Object {$_.Name -eq "
        "'ai-daily-office-parser.exe'}|ForEach-Object {[int]$_.ProcessId});"
        "$discovery=@($all|Where-Object {$_.Name -eq "
        "'ai-daily-discovery.exe'}|ForEach-Object {[int]$_.ProcessId});"
        "$document=@($all|Where-Object {($_.Name -eq 'python.exe' -or "
        "$_.Name -eq 'pythonw.exe') -and $_.CommandLine -like "
        "'*src.workers.document_parser_worker*'}|"
        "ForEach-Object {[int]$_.ProcessId});"
        "$legacy=@($all|Where-Object {($_.Name -eq 'python.exe' -or "
        "$_.Name -eq 'pythonw.exe') -and ($_.CommandLine -like "
        "'*multiprocessing.spawn*' -or $_.CommandLine -like "
        "'*spawn_main*')}|ForEach-Object {[int]$_.ProcessId});"
        "$pythonAll=@($all|Where-Object {$_.Name -eq 'python.exe' -or "
        "$_.Name -eq 'pythonw.exe'}|ForEach-Object {[int]$_.ProcessId});"
        "$rustTest=@($all|Where-Object {$_.ExecutablePath -like "
        "($env:AI_DAILY_GATE_PROJECT_ROOT+'\\rust\\target\\*')}|"
        "ForEach-Object {[int]$_.ProcessId});"
        "[pscustomobject]@{scanner=@($scanner);office=@($office);"
        "discovery=@($discovery);python_document_worker=@($document);"
        "legacy_spawn_worker=@($legacy);python_process=@($pythonAll);"
        "rust_target_process=@($rustTest)}|ConvertTo-Json -Compress"
    )
    try:
        completed = subprocess.run(
            [
                "powershell.exe",
                "-NoProfile",
                "-NonInteractive",
                "-Command",
                script,
            ],
            stdin=subprocess.DEVNULL,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=15.0,
            check=False,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
            env={
                **os.environ,
                "AI_DAILY_GATE_PROJECT_ROOT": str(PROJECT_ROOT),
            },
        )
        if completed.returncode != 0:
            raise RuntimeError("process query returned nonzero")
        raw = json.loads(completed.stdout.decode("utf-8", errors="strict"))
        snapshot = {
            key: _normalize_pid_list(raw[key])
            for key in (
                "scanner",
                "office",
                "discovery",
                "python_document_worker",
                "legacy_spawn_worker",
                "python_process",
                "rust_target_process",
            )
        }
    except (OSError, UnicodeError, ValueError, KeyError, TypeError) as exc:
        raise RuntimeError("scanner process snapshot failed") from exc
    return snapshot


def _wait_for_scanner_processes(
    baseline: dict[str, list[int]],
) -> tuple[dict[str, list[int]], bool]:
    snapshot = _scanner_process_snapshot()
    for _ in range(10):
        if not any(_new_process_pids(baseline, snapshot).values()):
            return snapshot, True
        time.sleep(0.2)
        snapshot = _scanner_process_snapshot()
    return snapshot, not any(_new_process_pids(baseline, snapshot).values())


def _process_delta_evidence(
    before: dict[str, list[int]],
    after: dict[str, list[int]],
) -> dict[str, Any]:
    new_pids = _new_process_pids(before, after)
    exited_pids = {
        key: sorted(set(before[key]) - set(after[key])) for key in before
    }
    return {
        "passed": not any(new_pids.values()),
        "before": dict(before),
        "after": dict(after),
        "new_pids": new_pids,
        "exited_pids": exited_pids,
    }


def _normalize_pid_list(value: Any) -> list[int]:
    values = value if isinstance(value, list) else [value]
    normalized = sorted({int(item) for item in values if item is not None})
    if any(item <= 0 for item in normalized):
        raise ValueError("process snapshot contains an invalid PID")
    return normalized


def _new_process_pids(
    before: dict[str, list[int]],
    after: dict[str, list[int]],
) -> dict[str, list[int]]:
    if set(before) != set(after):
        raise ValueError("process snapshots use different categories")
    return {
        key: sorted(set(after[key]) - set(before[key])) for key in before
    }


def all_cutover_gates_pass(*gates: dict[str, Any]) -> bool:
    return bool(gates) and all(gate.get("passed") is True for gate in gates)


@contextmanager
def _temporary_real_comparison_root() -> Iterator[Path]:
    temp_base = Path(tempfile.gettempdir()).resolve(strict=True)
    created = Path(
        tempfile.mkdtemp(prefix="ai-daily-real-compare-", dir=temp_base)
    ).resolve(strict=True)
    if created.parent != temp_base or not created.name.startswith(
        "ai-daily-real-compare-"
    ):
        raise RuntimeError("real comparison root failed safety validation")
    try:
        yield created
    finally:
        resolved = created.resolve(strict=True)
        if resolved.parent != temp_base or not resolved.name.startswith(
            "ai-daily-real-compare-"
        ):
            raise RuntimeError("refusing to remove unsafe real comparison root")
        gc.collect()
        last_error: OSError | None = None
        for _ in range(30):
            if not resolved.exists():
                last_error = None
                break
            try:
                shutil.rmtree(resolved)
                last_error = None
                break
            except OSError as exc:
                last_error = exc
                time.sleep(0.1)
                gc.collect()
        if last_error is not None:
            raise RuntimeError("real comparison cleanup failed") from last_error


def _percentile(values: list[float], percentile: float) -> float:
    ordered = sorted(float(value) for value in values)
    position = (len(ordered) - 1) * percentile
    lower = math.floor(position)
    upper = math.ceil(position)
    if lower == upper:
        return ordered[lower]
    fraction = position - lower
    return ordered[lower] + (ordered[upper] - ordered[lower]) * fraction


def _increment(counts: dict[str, int], key: str) -> None:
    counts[key] = counts.get(key, 0) + 1


def _require(violations: list[str], condition: bool, code: str) -> None:
    if not condition:
        violations.append(code)


def _new_directory(path: Path) -> Path:
    path.mkdir(parents=True)
    return path.resolve(strict=True)


def _parse_date(value: str) -> date:
    try:
        return date.fromisoformat(value)
    except ValueError as exc:
        raise argparse.ArgumentTypeError("date must use YYYY-MM-DD") from exc


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Run sanitized scanner parity/cache/fault/real/performance gates"
        ),
    )
    parser.add_argument("--artifact-root", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--scan-date", type=_parse_date, default=date.today())
    parser.add_argument("--samples", type=int, default=5)
    return parser


def main() -> int:
    args = build_parser().parse_args()
    logging.getLogger("ai_daily_report").setLevel(logging.CRITICAL + 1)
    try:
        root = args.artifact_root.resolve(strict=True)
        if not root.is_dir() or any(root.iterdir()):
            raise ValueError("artifact root must exist and be empty")
        output = args.output
        if not output.is_absolute():
            output = root / output
        output = output.resolve(strict=False)
        if not output.is_relative_to(root):
            raise ValueError("gate output must stay under artifact root")
        real_work_value = os.environ.get(_REAL_WORK_DIR_ENV, "").strip()
        real_sample_value = os.environ.get(_REAL_SAMPLE_FILE_ENV, "").strip()
        if not real_work_value:
            raise ValueError("AI_DAILY_REAL_WORK_DIR is required")
        if not real_sample_value:
            raise ValueError("AI_DAILY_REAL_SAMPLE_FILE is required")
        processes_before = _scanner_process_snapshot()
        try:
            gate_config = build_gate_config(app_config)
            real_gate_config = build_real_gate_config(app_config)
            parity = run_parity_gate(
                gate_root=_new_directory(root / "parity"),
                gate_config=gate_config,
                scan_date=args.scan_date,
            )
            cache = run_cache_gate(
                gate_root=_new_directory(root / "cache"),
                gate_config=gate_config,
                scan_date=args.scan_date,
            )
            performance = run_performance_gate(
                gate_root=_new_directory(root / "performance"),
                gate_config=gate_config,
                scan_date=args.scan_date,
                samples=args.samples,
            )
            fault = run_fault_gate(
                gate_root=_new_directory(root / "fault"),
            )
            real_directory = run_real_directory_gate(
                gate_root=_new_directory(root / "real-directory"),
                gate_config=real_gate_config,
                work_dir=Path(real_work_value),
                sample_file=Path(real_sample_value),
            )
        finally:
            processes_after, processes_settled = _wait_for_scanner_processes(
                processes_before,
            )
        process_gate = _process_delta_evidence(
            processes_before,
            processes_after,
        )
        process_gate["passed"] = (
            process_gate["passed"] and processes_settled
        )
        evidence = {
            "contract": "ai_daily_scanner_cutover_gate",
            "protocol_version": 1,
            "content_policy": "hashes_and_metadata_only",
            "passed": all_cutover_gates_pass(
                parity,
                cache,
                fault,
                performance,
                real_directory,
                process_gate,
            ),
            "parity": parity,
            "cache": cache,
            "fault": fault,
            "performance": performance,
            "real_directory": real_directory,
            "process_cleanup": process_gate,
        }
        comparison._assert_metadata_only(evidence)
        output.write_text(
            json.dumps(evidence, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    except Exception as exc:
        print(
            f"scanner cutover gate failed ({type(exc).__name__})",
            file=sys.stderr,
        )
        return 1
    print(
        json.dumps(
            {
                "passed": evidence["passed"],
                "parity_passed": parity["passed"],
                "cache_passed": cache["passed"],
                "fault_passed": fault["passed"],
                "performance_passed": performance["passed"],
                "real_directory_passed": real_directory["passed"],
                "process_cleanup_passed": process_gate["passed"],
                "performance_stats": performance["performance"]["stats"],
            },
            sort_keys=True,
        )
    )
    return 0 if evidence["passed"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
